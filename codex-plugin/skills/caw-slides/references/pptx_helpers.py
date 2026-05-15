"""Shared helpers for generating research presentation slides.

Encodes the caw-slides style guide (see ``references/style-guide.md``) so individual
``generate_<purpose>_<YYYYMMDD>.py`` scripts stay short and consistent.

All public functions return the created object so callers can further
tweak it when the style guide does not cover an edge case.

Layout: 16:9 (13.33" x 7.5"). Bilingual font split: 和文 = MS Gothic,
英数字 = Arial (auto-detected per character by ``mixed_runs`` / ``Run.font=None``).
Drop-in to any chemistry project — no project-specific names or paths.
"""
from __future__ import annotations

import os
from dataclasses import dataclass
from pathlib import Path

import matplotlib.font_manager as fm
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Fonts
# ---------------------------------------------------------------------------

FONT_JA = "MS Gothic"
FONT_EN = "Arial"

# Cross-platform MS Gothic candidates. First-existing path wins.
# Override with the ``CAW_SLIDES_MSGOTHIC`` environment variable if needed.
_MSGOTHIC_CANDIDATES: tuple[Path, ...] = (
    # macOS — bundled with Microsoft PowerPoint
    Path("/Applications/Microsoft PowerPoint.app/Contents/Resources/DFonts/msgothic.ttc"),
    # macOS — bundled with Microsoft Word (also valid)
    Path("/Applications/Microsoft Word.app/Contents/Resources/DFonts/msgothic.ttc"),
    # Windows
    Path("C:/Windows/Fonts/msgothic.ttc"),
    Path("C:/Windows/Fonts/MSGOTHIC.TTC"),
    # Linux — common install paths if user has set up MS Gothic manually
    Path("/usr/share/fonts/truetype/msgothic/msgothic.ttc"),
    Path("/usr/share/fonts/MS/msgothic.ttc"),
)


def _find_msgothic_path() -> Path | None:
    """Locate MS Gothic on disk. Returns ``None`` if it is unavailable.

    Honors the ``CAW_SLIDES_MSGOTHIC`` env var first so a user can point to a
    non-standard install (``export CAW_SLIDES_MSGOTHIC=/path/to/msgothic.ttc``).
    """
    env_override = os.environ.get("CAW_SLIDES_MSGOTHIC")
    if env_override:
        candidate = Path(env_override)
        if candidate.is_file():
            return candidate
    for candidate in _MSGOTHIC_CANDIDATES:
        if candidate.is_file():
            return candidate
    return None


# Resolved at import time. ``None`` on platforms without MS Gothic — callers
# that need Japanese in matplotlib output should detect this and either
# install MS Gothic or fall back to a CJK-capable system font.
MSGOTHIC_PATH = _find_msgothic_path()


def configure_matplotlib_japanese() -> str:
    """Register MS Gothic with matplotlib and set it as the default family.

    Returns
    -------
    str
        The matplotlib font-family name that was activated.

    Raises
    ------
    FileNotFoundError
        If MS Gothic is not found on any of the platform-default paths and
        ``CAW_SLIDES_MSGOTHIC`` is not set.

    Notes
    -----
    Must be called before any matplotlib rendering that includes Japanese
    text. Without this, characters render as tofu (縦長の □).
    """
    if MSGOTHIC_PATH is None:
        raise FileNotFoundError(
            "MS Gothic not found on any candidate path. Install Microsoft "
            "PowerPoint (macOS/Windows) or set CAW_SLIDES_MSGOTHIC to a "
            ".ttc/.ttf path."
        )
    fm.fontManager.addfont(str(MSGOTHIC_PATH))
    family = fm.FontProperties(fname=str(MSGOTHIC_PATH)).get_name()
    plt.rcParams["font.family"] = family
    plt.rcParams["axes.spines.top"] = False
    plt.rcParams["axes.spines.right"] = False
    return family


# ---------------------------------------------------------------------------
# Palette (see CLAUDE.md §3)
# ---------------------------------------------------------------------------

COLOR_TEXT_BODY = RGBColor(0x22, 0x22, 0x22)
COLOR_EMPH_BLUE = RGBColor(0x00, 0x70, 0xC0)
COLOR_EMPH_NAVY = RGBColor(0x00, 0x33, 0xCC)
COLOR_EMPH_RED = RGBColor(0xFF, 0x00, 0x00)
COLOR_SUB_CYAN = RGBColor(0x00, 0xB0, 0xF0)
COLOR_SUB_GREEN = RGBColor(0x00, 0xAC, 0x48)

COLOR_ACCENT_BLUE = RGBColor(0x44, 0x72, 0xC4)
COLOR_ACCENT_LIGHTBLUE = RGBColor(0x00, 0xAA, 0xFF)
COLOR_ACCENT_CYAN = RGBColor(0x2B, 0xBC, 0xE3)
COLOR_ACCENT_RED = RGBColor(0xFF, 0x50, 0x50)
COLOR_HIGHLIGHT_YELLOW = RGBColor(0xFF, 0xFF, 0x00)
COLOR_HIGHLIGHT_AMBER = RGBColor(0xFF, 0xC0, 0x00)

# 7原則のタイトル・強調色 (§12)
COLOR_TITLE = RGBColor(0x1A, 0x56, 0xA0)
COLOR_EMPH_7PRINCIPLES_RED = RGBColor(0xC0, 0x39, 0x3A)

# ---------------------------------------------------------------------------
# Layout constants (16:9 slide, §1)
# ---------------------------------------------------------------------------

SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

MARGIN_LEFT = Inches(0.4)
TITLE_TOP = Inches(0.12)
TITLE_HEIGHT = Inches(0.7)
SEPARATOR_Y = Inches(0.82)
BODY_TOP = Inches(0.95)
BODY_HEIGHT = Inches(5.6)
KEY_MESSAGE_Y = Inches(6.8)

# Font sizes (§2, §12)
SIZE_TITLE_MAIN = Pt(72)
SIZE_THEME = Pt(32)
SIZE_SLIDE_TITLE = Pt(28)
SIZE_SECTION = Pt(24)
SIZE_BODY = Pt(20)
SIZE_EMPH = Pt(24)
SIZE_NOTE = Pt(12)


# ---------------------------------------------------------------------------
# Presentation bootstrap
# ---------------------------------------------------------------------------


def new_presentation() -> Presentation:
    """Create a blank 16:9 presentation sized per the style guide."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def blank_slide(prs: Presentation):
    """Add a blank slide (layout 6 is the blank layout in default master)."""
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------------------------------------------------------------------------
# Text helpers
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class Run:
    """A single styled text run. Use multiple Run objects to mix JA/EN."""

    text: str
    size: Pt = SIZE_BODY
    bold: bool = False
    color: RGBColor = COLOR_TEXT_BODY
    # None = auto: JA if any Hiragana/Katakana/Kanji, else EN
    font: str | None = None


def _is_japanese(text: str) -> bool:
    for ch in text:
        code = ord(ch)
        if (
            0x3000 <= code <= 0x303F  # CJK symbols/punctuation (、 。 「 」 …)
            or 0x3040 <= code <= 0x309F  # Hiragana
            or 0x30A0 <= code <= 0x30FF  # Katakana
            or 0x4E00 <= code <= 0x9FFF  # CJK Unified Ideographs
            or 0xFF00 <= code <= 0xFFEF  # Halfwidth / Fullwidth forms
        ):
            return True
    return False


def add_text_box(
    slide,
    runs: list[Run],
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    word_wrap: bool = True,
    auto_size: bool = True,
):
    """Add a textbox whose runs are individually font-styled.

    Parameters
    ----------
    runs
        Sequence of :class:`Run`. Mix JA/EN runs to satisfy §1 (MS Gothic
        for 和文, Arial for 英数字) within a single textbox.
    auto_size
        If True, enables text-autofit-to-shape. **Do not rely on this to
        fit arbitrary amounts of text** — see §12 layout verification.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = anchor
    if auto_size:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    for idx, run_spec in enumerate(runs):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = run_spec.text
        font = run.font
        font.size = run_spec.size
        font.bold = run_spec.bold
        font.color.rgb = run_spec.color
        font.name = run_spec.font or (
            FONT_JA if _is_japanese(run_spec.text) else FONT_EN
        )

    return box


def add_title(slide, text: str, *, slide_number: int | None = None):
    """Add the standard slide title + blue separator line (§12 layout)."""
    add_text_box(
        slide,
        [Run(text, size=SIZE_SLIDE_TITLE, bold=True, color=COLOR_TITLE)],
        left=MARGIN_LEFT,
        top=TITLE_TOP,
        width=Inches(11.0),
        height=TITLE_HEIGHT,
        auto_size=False,
    )
    line = slide.shapes.add_connector(
        1, MARGIN_LEFT, SEPARATOR_Y, SLIDE_WIDTH - MARGIN_LEFT, SEPARATOR_Y
    )
    line.line.color.rgb = COLOR_ACCENT_BLUE
    line.line.width = Pt(1.5)

    if slide_number is not None:
        add_text_box(
            slide,
            [Run(str(slide_number), size=SIZE_NOTE, color=COLOR_TEXT_BODY)],
            left=SLIDE_WIDTH - Inches(0.8),
            top=TITLE_TOP,
            width=Inches(0.4),
            height=Inches(0.4),
            align=PP_ALIGN.RIGHT,
        )


def add_key_message(slide, runs: list[Run]):
    """Add the bottom key-message band (§12 layout)."""
    add_text_box(
        slide,
        runs,
        left=MARGIN_LEFT,
        top=KEY_MESSAGE_Y,
        width=SLIDE_WIDTH - 2 * MARGIN_LEFT,
        height=Inches(0.55),
        auto_size=False,
    )


# ---------------------------------------------------------------------------
# Image helpers
# ---------------------------------------------------------------------------


def add_picture_fit(
    slide,
    image_path: str | Path,
    *,
    left: Emu,
    top: Emu,
    max_width: Emu,
    max_height: Emu,
):
    """Add a picture preserving aspect ratio, centered in a bounding box.

    The **final** rendered width/height — not ``max_width``/``max_height``
    — is what you should reason about when checking for overlap with
    neighboring shapes (§12 layout verification).
    """
    from PIL import Image

    with Image.open(image_path) as img:
        iw, ih = img.size
    # Convert pixel dimensions to EMU assuming 96 DPI (1 px = 9525 EMU)
    px_w = Emu(int(iw * 9525))
    px_h = Emu(int(ih * 9525))
    ratio = min(max_width / px_w, max_height / px_h)
    target_w = Emu(int(iw * 9525 * ratio))
    target_h = Emu(int(ih * 9525 * ratio))
    x = left + (max_width - target_w) // 2
    y = top + (max_height - target_h) // 2
    return slide.shapes.add_picture(
        str(image_path), x, y, width=target_w, height=target_h
    )


# ---------------------------------------------------------------------------
# Forbidden-character guard (§11)
# ---------------------------------------------------------------------------

FORBIDDEN_IN_MATPLOTLIB = {
    "✓": "○",  # ✓
    "✗": "×",  # ✗
    "−": "-",  # Unicode minus
}
# Unicode subscript digits U+2080..U+2089 → ASCII
for _i in range(10):
    FORBIDDEN_IN_MATPLOTLIB[chr(0x2080 + _i)] = str(_i)


def sanitize_for_matplotlib(text: str) -> str:
    """Replace glyphs missing from MS Gothic with safe ASCII fallbacks."""
    for bad, good in FORBIDDEN_IN_MATPLOTLIB.items():
        text = text.replace(bad, good)
    return text


# ---------------------------------------------------------------------------
# Layout overlap guard (CLAUDE.md §「テキストボックス重なり禁止」)
# ---------------------------------------------------------------------------

Rect = tuple[int, int, int, int, str]


def assert_no_overlap(rects: list[Rect]) -> None:
    """Raise ``ValueError`` if any two rectangles in ``rects`` overlap.

    Each rectangle is ``(left, top, width, height, label)`` in EMU. Tangent
    edges (one rect ending exactly where the next begins) are **not** counted
    as overlapping -- honest adjacency is fine, visible intersection is not.
    """
    for i, a in enumerate(rects):
        ax, ay, aw, ah, al = a
        for b in rects[i + 1 :]:
            bx, by, bw, bh, bl = b
            if ax + aw <= bx or bx + bw <= ax:
                continue
            if ay + ah <= by or by + bh <= ay:
                continue
            raise ValueError(
                f"Layout overlap: {al!r} at ({ax},{ay},{aw}x{ah}) "
                f"overlaps {bl!r} at ({bx},{by},{bw}x{bh})"
            )


# ---------------------------------------------------------------------------
# Table helper (prefer tables over bullet lists for comparison data)
# ---------------------------------------------------------------------------


# ---------------------------------------------------------------------------
# Codex-style layout (2026-04-24: adopted from reference decks)
# ---------------------------------------------------------------------------

COLOR_CARD_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_CARD_BORDER = RGBColor(0xDD, 0xDD, 0xDD)
COLOR_PILL_BLUE_FILL = RGBColor(0xEA, 0xF2, 0xFC)
COLOR_PILL_RED_FILL = RGBColor(0xFB, 0xEA, 0xE8)
COLOR_PILL_GREEN_FILL = RGBColor(0xE8, 0xF6, 0xEE)
COLOR_PILL_GREY_FILL = RGBColor(0xF4, 0xF6, 0xF8)
COLOR_KEY_MSG_FILL = RGBColor(0xEA, 0xF2, 0xFC)
# Source-line colour is now quite light and the font is small — the source
# line must read as marginalia, not content.
COLOR_SOURCE_GREY = RGBColor(0xB0, 0xB0, 0xB0)
COLOR_LABEL_GREY = RGBColor(0x88, 0x88, 0x88)   # small caption labels (≤ 12pt)

SIZE_CARD_SECTION = Pt(21)  # ▸ headings inside body card
SIZE_CARD_BULLET = Pt(16)   # ・ bullet text
SIZE_PILL_TITLE = Pt(14)
SIZE_PILL_NUMBER = Pt(19)
SIZE_PILL_CAPTION = Pt(12)

# Fixed reference positions (copied from Codex reference decks)
CODEX_TITLE_RECT = (Inches(0.4), Inches(0.12), Inches(11.7), Inches(0.45))
CODEX_SLIDE_NUM_RECT = (Inches(12.18), Inches(0.12), Inches(0.62), Inches(0.4))
CODEX_SEP_RECT = (Inches(0.4), Inches(0.82), Inches(12.53), Inches(0.03))
CODEX_KEY_MSG_RECT = (Inches(0.52), Inches(6.28), Inches(12.25), Inches(0.58))
CODEX_SOURCE_RECT = (Inches(0.52), Inches(7.02), Inches(12.2), Inches(0.35))


@dataclass(frozen=True)
class Paragraph:
    """A single text paragraph composed of multiple styled runs.

    Use this when a paragraph needs to mix Japanese (MS Gothic) and English
    (Arial) runs, or to emphasize inline keywords in a different color.
    """

    runs: list["Run"]
    alignment: PP_ALIGN = PP_ALIGN.LEFT


def _apply_runs_to_paragraph(paragraph, runs: list["Run"]) -> None:
    for run_spec in runs:
        r = paragraph.add_run()
        r.text = run_spec.text
        r.font.size = run_spec.size
        r.font.bold = run_spec.bold
        r.font.color.rgb = run_spec.color
        r.font.name = run_spec.font or (
            FONT_JA if _is_japanese(run_spec.text) else FONT_EN
        )


def _write_paragraphs(text_frame, paragraphs: list[Paragraph]) -> None:
    for p_idx, para in enumerate(paragraphs):
        if p_idx == 0:
            p = text_frame.paragraphs[0]
            # Clear any default run
            for run in list(p.runs):
                run.text = ""
        else:
            p = text_frame.add_paragraph()
        p.alignment = para.alignment
        _apply_runs_to_paragraph(p, para.runs)


def add_rich_text_box(
    slide,
    paragraphs: list[Paragraph],
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    word_wrap: bool = True,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    auto_size: bool = False,
):
    """Add a text box supporting multiple paragraphs, each with multiple runs.

    Use this in place of :func:`add_text_box` whenever a paragraph mixes
    JA/EN or has inline color/weight changes.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = anchor
    if auto_size:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    _write_paragraphs(tf, paragraphs)
    return box


def _style_shape_fill(shape, fill: RGBColor | None, border: RGBColor | None) -> None:
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if border is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border
        shape.line.width = Pt(0.75)


def _set_shape_shadow(shape, on: bool) -> None:
    """Inject drop-shadow XML or suppress the theme-default shadow.

    PowerPoint themes apply a subtle drop-shadow to rectangles by default.
    Codex-style slides want a **flat** look everywhere except for a handful
    of emphasis shapes. Inject an empty ``<a:effectLst/>`` to override the
    theme shadow with "no effect", or an explicit ``<a:outerShdw>`` for an
    emphasis shape.
    """
    from lxml import etree
    from pptx.oxml.ns import qn

    spPr = shape._element.spPr
    existing = spPr.find(qn("a:effectLst"))
    if existing is not None:
        spPr.remove(existing)
    effectLst = etree.SubElement(spPr, qn("a:effectLst"))
    if on:
        shdw = etree.SubElement(
            effectLst, qn("a:outerShdw"),
            attrib={
                "blurRad": "38100",  # 0.1 cm
                "dist": "25400",     # 0.07 cm
                "dir": "2700000",    # 45 deg
                "algn": "tl",
                "rotWithShape": "0",
            },
        )
        col = etree.SubElement(shdw, qn("a:srgbClr"), attrib={"val": "000000"})
        etree.SubElement(col, qn("a:alpha"), attrib={"val": "30000"})


def mixed_runs(
    text: str,
    *,
    size: Pt,
    bold: bool = False,
    color: RGBColor = COLOR_TEXT_BODY,
) -> list["Run"]:
    """Split a mixed-language string into per-language runs.

    Each character is classified as Japanese (MS Gothic) or English
    (Arial, which covers ASCII, digits, Latin-extended). Consecutive
    characters of the same class are merged into one Run so paragraph
    formatting stays clean.
    """
    if not text:
        return []
    parts: list[tuple[str, bool]] = []
    cur_text = text[0]
    cur_ja = _is_japanese(text[0])
    for ch in text[1:]:
        is_ja = _is_japanese(ch)
        if is_ja == cur_ja:
            cur_text += ch
        else:
            parts.append((cur_text, cur_ja))
            cur_text = ch
            cur_ja = is_ja
    parts.append((cur_text, cur_ja))
    return [
        Run(segment, size=size, bold=bold, color=color,
            font=FONT_JA if is_ja else FONT_EN)
        for segment, is_ja in parts
    ]


def add_shape_card(
    slide,
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    paragraphs: list[Paragraph] | None = None,
    fill: RGBColor | None = None,
    border: RGBColor | None = None,
    padding: Emu = Inches(0.2),
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    shadow: bool = False,
):
    """Add a rectangular card (AUTO_SHAPE) optionally filled with paragraphs.

    The card is a single shape — its text lives inside the card's text
    frame. Unlike the old ``add_text_box`` + decorative background split,
    this keeps text and panel as one shape so ``assert_no_overlap`` only
    sees one rectangle per card.

    ``padding`` is applied as internal margins so bullet text does not kiss
    the rounded edge of the card.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    _style_shape_fill(shape, fill, border)
    _set_shape_shadow(shape, shadow)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = padding
    tf.margin_right = padding
    tf.margin_top = padding
    tf.margin_bottom = padding
    if paragraphs:
        _write_paragraphs(tf, paragraphs)
    return shape


def add_pill(
    slide,
    paragraphs: list[Paragraph],
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    fill: RGBColor = COLOR_PILL_BLUE_FILL,
    border: RGBColor | None = None,
    padding: Emu = Inches(0.06),
    shadow: bool = False,
):
    """Small colored callout pill (AUTO_SHAPE). Text is centered by default.

    Set ``shadow=True`` only for the one or two pills you want to emphasize
    on a slide (e.g. the headline number). Default is a flat Codex-style
    pill with no drop-shadow.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    _style_shape_fill(shape, fill, border)
    _set_shape_shadow(shape, shadow)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = padding
    tf.margin_right = padding
    tf.margin_top = padding
    tf.margin_bottom = padding
    _write_paragraphs(tf, paragraphs)
    return shape


def add_key_message_band(
    slide,
    paragraphs: list[Paragraph],
    *,
    fill: RGBColor = COLOR_KEY_MSG_FILL,
    border: RGBColor | None = RGBColor(0x22, 0x22, 0x22),
    shadow: bool = True,
):
    """Full-width key-message band at the Codex fixed position (y=6.28).

    Shadow defaults to ON because this band is the slide's headline
    takeaway — the one place where drop-shadow emphasis is always welcome.
    """
    left, top, width, height = CODEX_KEY_MSG_RECT
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    _style_shape_fill(shape, fill, border)
    _set_shape_shadow(shape, shadow)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    _write_paragraphs(tf, paragraphs)
    return shape


def add_source_line(slide, text: str):
    """Tiny source credit at the Codex fixed position (y=7.02).

    Pt 9 + light grey (#B0B0B0) — marginalia only, must not compete with
    the slide's content for the reader's attention.
    """
    left, top, width, height = CODEX_SOURCE_RECT
    runs = [Run(text, size=Pt(9), bold=False, color=COLOR_SOURCE_GREY,
                font=FONT_EN)]
    return add_rich_text_box(
        slide, [Paragraph(runs)], left=left, top=top, width=width,
        height=height, anchor=MSO_ANCHOR.TOP,
    )


def add_category_pill(
    slide,
    text: str,
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    fill: RGBColor = COLOR_PILL_BLUE_FILL,
    border: RGBColor | None = None,
    shadow: bool = False,
):
    """Title-slide category pill (e.g. '報告会 No.12')."""
    runs = mixed_runs(text, size=Pt(20), bold=True, color=COLOR_TITLE)
    return add_pill(
        slide,
        [Paragraph(runs, alignment=PP_ALIGN.CENTER)],
        left=left, top=top, width=width, height=height,
        fill=fill, border=border, shadow=shadow,
    )


# ---------------------------------------------------------------------------
# Native flow diagram (editable PowerPoint shapes, slides 3 / 9)
# ---------------------------------------------------------------------------


def add_flow_box(
    slide,
    text: str,
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    fill: RGBColor,
    text_color: RGBColor = COLOR_TEXT_BODY,
    font_size: Pt = Pt(12),
    bold: bool = True,
    shadow: bool = False,
    shape_type: MSO_SHAPE = MSO_SHAPE.ROUNDED_RECTANGLE,
):
    """Single box for a flow/timeline diagram with centered multi-line text.

    ``text`` may contain ``\\n`` for line breaks; each line becomes its own
    paragraph so PowerPoint wraps cleanly. Language-aware font switching is
    applied per paragraph.
    """
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    _style_shape_fill(shape, fill, None)
    _set_shape_shadow(shape, shadow)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    paragraphs = [
        Paragraph(
            mixed_runs(line, size=font_size, bold=bold, color=text_color),
            alignment=PP_ALIGN.CENTER,
        )
        for line in text.split("\n")
    ]
    _write_paragraphs(tf, paragraphs)
    return shape


def add_flow_arrow(
    slide,
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    direction: str = "right",
    fill: RGBColor = COLOR_TEXT_BODY,
    border: RGBColor | None = None,
):
    """Add a filled arrow shape (``MSO_SHAPE.<direction>_ARROW``).

    A dedicated arrow auto-shape is used instead of a ``MSO_CONNECTOR``
    because connectors with zero-width/zero-height bounding boxes (purely
    horizontal or vertical) occasionally serialize coordinates as floats
    that break python-pptx round-trips. A filled auto-shape has an
    explicit rectangular bounding box and is trivially selectable /
    editable in PowerPoint.

    ``direction`` is one of ``"right"``, ``"left"``, ``"up"``, ``"down"``.
    """
    shape_map = {
        "right": MSO_SHAPE.RIGHT_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
    }
    shape_type = shape_map[direction]
    shape = slide.shapes.add_shape(
        shape_type,
        int(left), int(top), int(width), int(height),
    )
    _style_shape_fill(shape, fill, border)
    _set_shape_shadow(shape, False)
    # Arrows are decorative; do not add text.
    return shape


# ---------------------------------------------------------------------------
# Native Excel-editable charts (replaces matplotlib PNGs for graphs)
# ---------------------------------------------------------------------------


def _style_chart_common(
    chart,
    *,
    title: str | None,
    show_legend: bool = True,
    legend_font: Pt = Pt(11),
    title_font: Pt = Pt(13),
    axis_font: Pt = Pt(10),
    plot_area_border: RGBColor | None = RGBColor(0x00, 0x00, 0x00),
    plot_area_border_width: Pt = Pt(1.0),
    gridlines: bool = False,
    gridlines_color: RGBColor = RGBColor(0xDD, 0xDD, 0xDD),
    gridlines_width: Pt = Pt(0.5),
) -> None:
    """Apply consistent styling across chart types.

    Default behavior (per 2026-04-24 feedback):
    - **Plot-area border**: black 1pt (clear publication-style frame).
    - **Gridlines**: OFF. Data readability takes priority over reference
      lines. Pass ``gridlines=True`` to enable subtle thin pale-grey
      reference lines when a specific chart needs them.
    """
    if title:
        chart.has_title = True
        tframe = chart.chart_title.text_frame
        tframe.text = title
        for p in tframe.paragraphs:
            for r in p.runs:
                r.font.size = title_font
                r.font.bold = True
                r.font.name = FONT_JA if _is_japanese(r.text) else FONT_EN
    else:
        chart.has_title = False
    chart.has_legend = show_legend
    if show_legend:
        try:
            chart.legend.include_in_layout = False
            legend_font_obj = chart.legend.font
            legend_font_obj.size = legend_font
            legend_font_obj.name = FONT_EN
        except Exception:
            pass
    try:
        for axis in (chart.category_axis, chart.value_axis):
            if axis is None:
                continue
            axis.tick_labels.font.size = axis_font
            axis.tick_labels.font.name = FONT_EN
    except Exception:
        pass

    # Plot-area border: python-pptx 1.0.2 doesn't expose Chart.plot_area,
    # so we inject <c:spPr>/<a:ln> into <c:plotArea> via XML. The spPr must
    # be the last child of plotArea (just before an optional c:extLst) per
    # the OOXML schema.
    if plot_area_border is not None:
        try:
            from lxml import etree
            from pptx.oxml.ns import qn
            chart_space = chart.element  # <c:chartSpace>
            chart_xml = chart_space.find(qn("c:chart"))
            plot_area_xml = chart_xml.find(qn("c:plotArea"))
            # Remove any existing spPr so we replace cleanly
            for existing in plot_area_xml.findall(qn("c:spPr")):
                plot_area_xml.remove(existing)
            # Detach extLst if present, append new spPr, then re-append extLst
            extLst = plot_area_xml.find(qn("c:extLst"))
            if extLst is not None:
                plot_area_xml.remove(extLst)
            spPr = etree.SubElement(plot_area_xml, qn("c:spPr"))
            # Transparent plot-area fill so the data region reads as "paper"
            etree.SubElement(spPr, qn("a:noFill"))
            # a:ln w="..." expects EMU. pptx.util.Pt is already an int in
            # EMU (1 pt = 12700 EMU), so int(Pt(1.0)) == 12700 is correct.
            ln = etree.SubElement(spPr, qn("a:ln"))
            ln.set("w", str(int(plot_area_border_width)))
            solidFill = etree.SubElement(ln, qn("a:solidFill"))
            clr = etree.SubElement(solidFill, qn("a:srgbClr"))
            clr.set("val", f"{plot_area_border}")
            if extLst is not None:
                plot_area_xml.append(extLst)
        except Exception:
            pass

    # Gridlines
    for axis in (chart.value_axis, chart.category_axis):
        if axis is None:
            continue
        try:
            axis.has_major_gridlines = bool(gridlines)
            if gridlines:
                gl = axis.major_gridlines
                gline = gl.format.line
                gline.color.rgb = gridlines_color
                gline.width = gridlines_width
        except Exception:
            pass


def add_bar_chart(
    slide,
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    categories: list[str],
    series_data: list[tuple[str, list[float]]],
    title: str | None = None,
    y_label: str | None = None,
    gridlines: bool = False,
):
    """Excel-editable clustered-column bar chart.

    ``series_data`` is ``[(series_name, values), ...]``. The resulting chart
    shape opens an embedded workbook on double-click so the user can edit
    numbers after the deck is generated.
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_data:
        chart_data.add_series(name, values)
    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    )
    chart = frame.chart
    _style_chart_common(chart, title=title, gridlines=gridlines)
    if y_label is not None:
        try:
            chart.value_axis.has_title = True
            chart.value_axis.axis_title.text_frame.text = y_label
        except Exception:
            pass
    return frame


def add_scatter_line_chart(
    slide,
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    series_data: list[tuple[str, list[float], list[float]]],
    title: str | None = None,
    x_label: str | None = None,
    y_label: str | None = None,
    show_markers: bool = True,
    show_legend: bool = True,
    gridlines: bool = False,
):
    """Excel-editable XY scatter-line chart.

    ``series_data`` is ``[(series_name, x_values, y_values), ...]``.
    """
    from pptx.chart.data import XyChartData
    from pptx.enum.chart import XL_CHART_TYPE

    chart_data = XyChartData()
    for name, xs, ys in series_data:
        s = chart_data.add_series(name)
        for x, y in zip(xs, ys):
            s.add_data_point(x, y)
    chart_type = (
        XL_CHART_TYPE.XY_SCATTER_LINES
        if show_markers
        else XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS
    )
    frame = slide.shapes.add_chart(
        chart_type, left, top, width, height, chart_data
    )
    chart = frame.chart
    _style_chart_common(chart, title=title, show_legend=show_legend,
                        gridlines=gridlines)
    if x_label is not None:
        try:
            chart.category_axis.has_title = True
            chart.category_axis.axis_title.text_frame.text = x_label
        except Exception:
            pass
    if y_label is not None:
        try:
            chart.value_axis.has_title = True
            chart.value_axis.axis_title.text_frame.text = y_label
        except Exception:
            pass
    return frame


def add_slide_chrome(slide, title: str, slide_number: int) -> list[Rect]:
    """Add Codex slide chrome: title text, slide number, blue separator bar.

    Returns a list of non-overlapping rects (title, slide number, separator)
    that callers should include in their ``assert_no_overlap`` check.
    """
    # Title text
    t_left, t_top, t_w, t_h = CODEX_TITLE_RECT
    add_rich_text_box(
        slide,
        [Paragraph([Run(title, size=Pt(28), bold=True, color=COLOR_TITLE)])],
        left=t_left, top=t_top, width=t_w, height=t_h,
    )
    # Slide number
    n_left, n_top, n_w, n_h = CODEX_SLIDE_NUM_RECT
    add_rich_text_box(
        slide,
        [Paragraph(
            [Run(str(slide_number), size=Pt(12), bold=True,
                 color=COLOR_SOURCE_GREY, font=FONT_EN)],
            alignment=PP_ALIGN.RIGHT,
        )],
        left=n_left, top=n_top, width=n_w, height=n_h,
    )
    # Blue separator bar (AUTO_SHAPE so it's visible as a filled strip)
    s_left, s_top, s_w, s_h = CODEX_SEP_RECT
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, s_left, s_top, s_w, s_h
    )
    _style_shape_fill(bar, COLOR_ACCENT_BLUE, None)
    return [
        (t_left, t_top, t_w, t_h, "<title>"),
        (n_left, n_top, n_w, n_h, "<slide-number>"),
        (s_left, s_top, s_w, s_h, "<separator>"),
    ]


# ---------------------------------------------------------------------------
# Legacy data-table helper
# ---------------------------------------------------------------------------


def add_data_table(
    slide,
    headers: list[str],
    rows: list[list[str]],
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    header_fill: RGBColor = COLOR_ACCENT_BLUE,
    header_text: RGBColor = RGBColor(0xFF, 0xFF, 0xFF),
    body_text: RGBColor = COLOR_TEXT_BODY,
    font_size: Pt = Pt(16),
    header_size: Pt = Pt(16),
    first_col_bold: bool = True,
):
    """Add an editable native PowerPoint table.

    Font is auto-selected per cell (MS Gothic for JA, Arial for EN) to honor
    section 1 of the style guide. Header row is filled with ``header_fill``.
    """
    n_rows = len(rows) + 1
    n_cols = len(headers)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = shape.table

    for c, htext in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.text_frame.word_wrap = True
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = htext
        run.font.size = header_size
        run.font.bold = True
        run.font.color.rgb = header_text
        run.font.name = FONT_JA if _is_japanese(htext) else FONT_EN

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text_frame.word_wrap = True
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = value
            run.font.size = font_size
            run.font.bold = first_col_bold and c == 0
            run.font.color.rgb = body_text
            run.font.name = FONT_JA if _is_japanese(value) else FONT_EN

    return shape
