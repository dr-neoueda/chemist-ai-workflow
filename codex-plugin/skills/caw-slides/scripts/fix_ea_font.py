#!/usr/bin/env python3
"""ビルド後 pptx の East-Asian フォント修正（caw-slides）。

SVG→pptx 変換器は、MS Gothic 親 text 内の Arial ``<tspan>`` の直後に続く「かな tail
テキスト」を、直前 tspan の ea=Arial を継いだ run として書き出す癖がある。すると
PowerPoint が CJK 文字をシステム CJK フォントに代替し MS Gothic と不一致になる。
この後処理は **CJK コードポイントを含む全 run の ``<a:ea>`` typeface を日本語フォント
（既定 MS Gothic）に強制**する（``<a:latin>`` は据え置き）。

assert_font_rule（SVG 段の authoring ミス検出）とは別レイヤ：こちらは変換器の癖に
対する pptx 段の機械的修正。フォントゲート緑でも本処理は必ず 1 回かける。

Usage
-----
    python3 fix_ea_font.py <pptx>            # 既定 MS Gothic に修正、上書き保存
    python3 fix_ea_font.py <pptx> --font "Yu Gothic"

判定ロジックのテストは ``tests/test_fix_ea_font.py`` (pytest) を参照。
依存: python-pptx。
"""
from __future__ import annotations

import argparse
import os
import sys
from typing import Any, Iterator

# East-Asian フォントを要する文字範囲（かな・漢字・全角・囲み CJK ①②③ 等）
_CJK_RANGES: tuple[tuple[int, int], ...] = (
    (0x3000, 0x30FF),   # 句読点・ひらがな・カタカナ・全角記号
    (0x3200, 0x33FF),   # 囲み CJK（丸数字 ①②③ 等）・CJK 互換
    (0x3400, 0x9FFF),   # CJK 統合漢字（拡張 A 含む）
    (0xF900, 0xFAFF),   # CJK 互換漢字
    (0xFF00, 0xFFEF),   # 半角・全角形
)

DEFAULT_JA_FONT = "MS Gothic"


def has_cjk(text: str) -> bool:
    """text に East-Asian フォントを要する文字が含まれるか。"""
    return any(any(lo <= ord(ch) <= hi for lo, hi in _CJK_RANGES) for ch in text)


def _iter_shapes(shapes: Any) -> Iterator[Any]:
    """グループを再帰的に展開して全 shape を yield する。"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for sh in shapes:
        yield sh
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(sh.shapes)


def _iter_runs(shapes: Any) -> Iterator[Any]:
    """text_frame / native table cell を持つ shape の全 run を yield する。"""
    for sh in _iter_shapes(shapes):
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                yield from para.runs
        elif getattr(sh, "has_table", False):
            for row in sh.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        yield from para.runs


# CT_TextCharacterProperties の子要素順（<a:ea> は <a:latin> の後・<a:cs> の前）
_EA_SUCCESSORS = ("a:cs", "a:sym", "a:hlinkClick", "a:hlinkMouseOver", "a:rtl", "a:extLst")


def get_or_add_ea(rPr: Any) -> Any:
    """``<a:rPr>`` に ``<a:ea>``（East-Asian フォント）子要素を取得/追加する。

    python-pptx は ``ea`` を生成子として定義しないため、スキーマ順を守って lxml で挿入する。
    ``insert_element_before`` は内部で ``qn()`` を適用するため、後続要素は接頭辞記法
    （``"a:cs"`` 等）で渡す（Clark 記法で渡すと二重適用になり誤る）。挿入順の正しさは
    ``tests/test_fix_ea_font.py::test_ea_inserted_before_cs`` で検証。
    """
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls, qn

    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = parse_xml(f'<a:ea {nsdecls("a")}/>')
        rPr.insert_element_before(ea, *_EA_SUCCESSORS)
    return ea


def fix_ea_font(pptx_path: str, ja_font: str = DEFAULT_JA_FONT) -> int:
    """pptx を開き、CJK を含む run の ``<a:ea>`` を ja_font に強制する。

    Returns
    -------
    int
        typeface を書き換えた run 数（0 なら保存しない）。
    """
    from pptx import Presentation

    prs = Presentation(pptx_path)
    fixed = 0
    for slide in prs.slides:
        for run in _iter_runs(slide.shapes):
            if not has_cjk(run.text):
                continue
            rPr = run._r.get_or_add_rPr()
            ea = get_or_add_ea(rPr)
            if ea.get("typeface") != ja_font:
                ea.set("typeface", ja_font)
                fixed += 1
    if fixed:
        # アトミック保存（保存中断で元 pptx を壊さない）
        tmp = f"{pptx_path}.tmp"
        prs.save(tmp)
        os.replace(tmp, pptx_path)
    return fixed


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="ビルド後 pptx の ea フォント修正 (fix_ea_font)")
    parser.add_argument("pptx", help="対象 .pptx（上書き保存）")
    parser.add_argument("--font", default=DEFAULT_JA_FONT, help=f"日本語フォント名（既定 {DEFAULT_JA_FONT}）")
    args = parser.parse_args(argv)
    try:
        fixed = fix_ea_font(args.pptx, args.font)
    except Exception as exc:  # ファイル無し・不正 pptx 等は CLI 境界で握る
        print(f"[fix_ea_font] ERROR: {exc}", file=sys.stderr)
        return 1
    print(f"[fix_ea_font] CJK run の ea を {args.font!r} に修正: {fixed} 箇所")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
