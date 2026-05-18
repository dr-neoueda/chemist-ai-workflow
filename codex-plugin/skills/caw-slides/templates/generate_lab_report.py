"""Lab report / weekly progress template (6–15 slides).

caw-slides variant: **研究室報告会・進捗共有**

Use case
--------
In-lab weekly or biweekly progress meeting. Self-generated data is primary
(charts, tables, schemes), driven from analyses you ran this week. Less
formal than a conference deck but still follows the same style rules.

Structure
---------
1. Title slide (報告会 No. XX)
2. Theme reminder (1 slide)
3. Experimental / computational progress (2–8 slides)
4. Analysis / interpretation (1–3 slides)
5. Next steps (1 slide)

Usage
-----
1. Copy this file to ``.company/presentation/scripts/generate_report_<YYYYMMDD>.py``
2. Replace placeholders + plug in your data
3. Run — output to ``presentations/slides/report_<YYYYMMDD>.pptx``
"""
from __future__ import annotations

import sys
from pathlib import Path

# Locate ``pptx_helpers.py``: script-dir first (after user copies it alongside),
# then ../references/ (running this script directly from the plugin install).
_HERE = Path(__file__).resolve().parent
for _cand in (_HERE, _HERE.parent / "references"):
    if (_cand / "pptx_helpers.py").is_file():
        _cand_str = str(_cand)
        if _cand_str not in sys.path:
            sys.path.insert(0, _cand_str)
        break

import pptx_helpers as h  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

OUT_DIR = Path.cwd() / "presentations" / "slides"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_PATH = OUT_DIR / "lab_report_template.pptx"

# ─── Placeholders ─────────────────────────────────────────────────────────────
REPORT_NUMBER = "<XX>"
REPORT_DATE = "<YYYY-MM-DD>"
PRESENTER = "<your_name>"
LAB = "<lab name>"
THEME = "<研究テーマ>"


def build_title_slide(prs):
    slide = h.blank_slide(prs)
    rects = []

    h.add_rich_text_box(
        slide,
        [h.Paragraph(h.mixed_runs(f"報告会 No.{REPORT_NUMBER}",
                                  size=Pt(72), bold=True, color=h.COLOR_TITLE),
                     alignment=PP_ALIGN.CENTER)],
        left=Inches(1.0), top=Inches(2.2), width=Inches(11.3), height=Inches(1.6),
    )
    rects.append((Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.6), "<no>"))

    h.add_rich_text_box(
        slide,
        [h.Paragraph(h.mixed_runs(f"{REPORT_DATE}    {LAB}    {PRESENTER}",
                                  size=Pt(32), color=h.COLOR_TEXT_BODY),
                     alignment=PP_ALIGN.CENTER)],
        left=Inches(1.0), top=Inches(4.4), width=Inches(11.3), height=Inches(0.8),
    )
    rects.append((Inches(1.0), Inches(4.4), Inches(11.3), Inches(0.8), "<meta>"))
    h.assert_no_overlap(rects)


def build_theme_slide(prs, slide_number: int):
    slide = h.blank_slide(prs)
    rects = list(h.add_slide_chrome(slide, "研究テーマ", slide_number))

    h.add_shape_card(
        slide,
        left=Inches(0.4), top=Inches(1.18), width=Inches(12.53), height=Inches(4.95),
        paragraphs=[
            h.Paragraph(h.mixed_runs(THEME, size=Pt(28), bold=True, color=h.COLOR_TITLE),
                        alignment=PP_ALIGN.CENTER),
            h.Paragraph(h.mixed_runs("", size=Pt(16))),
            h.Paragraph(h.mixed_runs("▸ 目的", size=Pt(21), bold=True, color=h.COLOR_TITLE)),
            h.Paragraph(h.mixed_runs("・<目的の一文>", size=Pt(16))),
            h.Paragraph(h.mixed_runs("", size=Pt(16))),
            h.Paragraph(h.mixed_runs("▸ 今回の報告範囲",
                                     size=Pt(21), bold=True, color=h.COLOR_TITLE)),
            h.Paragraph(h.mixed_runs("・<報告範囲>", size=Pt(16))),
        ],
    )
    rects.append((Inches(0.4), Inches(1.18), Inches(12.53), Inches(4.95), "<theme>"))

    h.add_key_message_band(
        slide,
        [h.Paragraph(h.mixed_runs("<L1: 今週の主要トピック>",
                                  size=Pt(20), bold=True, color=h.COLOR_TEXT_BODY),
                     alignment=PP_ALIGN.CENTER)],
    )
    rects.append(h.CODEX_KEY_MSG_RECT + ("<key-msg>",))
    h.assert_no_overlap(rects)


def build_progress_chart_slide(prs, slide_number: int):
    """Progress slide with one native chart + L1 takeaway."""
    slide = h.blank_slide(prs)
    rects = list(h.add_slide_chrome(slide, "<今週の進捗: 実験／解析データ>", slide_number))

    h.add_scatter_line_chart(
        slide,
        left=Inches(0.4), top=Inches(1.18), width=Inches(12.53), height=Inches(4.6),
        series_data=[
            ("<run 1>", [0, 1, 2, 3, 4, 5], [0.2, 0.5, 0.9, 1.3, 1.6, 1.8]),
            ("<run 2>", [0, 1, 2, 3, 4, 5], [0.3, 0.6, 1.0, 1.2, 1.3, 1.3]),
        ],
        title=None, x_label="<x-axis label>", y_label="<y-axis label>",
        gridlines=False,
    )
    rects.append((Inches(0.4), Inches(1.18), Inches(12.53), Inches(4.6), "<chart>"))

    h.add_key_message_band(
        slide,
        [h.Paragraph(h.mixed_runs("<L1: データから言える一行>",
                                  size=Pt(20), bold=True, color=h.COLOR_TEXT_BODY),
                     alignment=PP_ALIGN.CENTER)],
    )
    rects.append(h.CODEX_KEY_MSG_RECT + ("<key-msg>",))
    h.assert_no_overlap(rects)


def build_progress_table_slide(prs, slide_number: int):
    """Progress slide with a comparison table (e.g. multiple runs)."""
    slide = h.blank_slide(prs)
    rects = list(h.add_slide_chrome(slide, "<条件比較: 実験／計算ジョブ>", slide_number))

    h.add_data_table(
        slide,
        headers=["<条件>", "<値 1>", "<値 2>", "<備考>"],
        rows=[
            ["<条件 A>", "<v1>", "<v2>", "<備考 A>"],
            ["<条件 B>", "<v1>", "<v2>", "<備考 B>"],
            ["<条件 C>", "<v1>", "<v2>", "<備考 C>"],
        ],
        left=Inches(0.4), top=Inches(1.18), width=Inches(12.53), height=Inches(4.6),
    )
    rects.append((Inches(0.4), Inches(1.18), Inches(12.53), Inches(4.6), "<table>"))

    h.add_key_message_band(
        slide,
        [h.Paragraph(h.mixed_runs("<L1: どの条件が良かったか・なぜか>",
                                  size=Pt(20), bold=True, color=h.COLOR_TEXT_BODY),
                     alignment=PP_ALIGN.CENTER)],
    )
    rects.append(h.CODEX_KEY_MSG_RECT + ("<key-msg>",))
    h.assert_no_overlap(rects)


def build_next_steps_slide(prs, slide_number: int):
    slide = h.blank_slide(prs)
    rects = list(h.add_slide_chrome(slide, "今後の予定", slide_number))

    h.add_shape_card(
        slide,
        left=Inches(0.4), top=Inches(1.18), width=Inches(12.53), height=Inches(4.95),
        paragraphs=[
            h.Paragraph(h.mixed_runs("▸ 来週やること",
                                     size=Pt(21), bold=True, color=h.COLOR_TITLE)),
            h.Paragraph(h.mixed_runs("・<タスク 1>", size=Pt(16))),
            h.Paragraph(h.mixed_runs("・<タスク 2>", size=Pt(16))),
            h.Paragraph(h.mixed_runs("", size=Pt(16))),
            h.Paragraph(h.mixed_runs("▸ 検討事項・相談",
                                     size=Pt(21), bold=True, color=h.COLOR_TITLE)),
            h.Paragraph(h.mixed_runs("・<相談事項 1>", size=Pt(16))),
        ],
    )
    rects.append((Inches(0.4), Inches(1.18), Inches(12.53), Inches(4.95), "<next>"))

    h.add_key_message_band(
        slide,
        [h.Paragraph(h.mixed_runs("<L1: 次の進捗報告で示したい結果>",
                                  size=Pt(20), bold=True, color=h.COLOR_TEXT_BODY),
                     alignment=PP_ALIGN.CENTER)],
    )
    rects.append(h.CODEX_KEY_MSG_RECT + ("<key-msg>",))
    h.assert_no_overlap(rects)


def main() -> None:
    prs = h.new_presentation()
    build_title_slide(prs)
    build_theme_slide(prs, slide_number=2)
    build_progress_chart_slide(prs, slide_number=3)
    build_progress_table_slide(prs, slide_number=4)
    build_next_steps_slide(prs, slide_number=5)
    prs.save(str(OUT_PATH))
    print(f"Saved {OUT_PATH}")


if __name__ == "__main__":
    main()
