"""Showcase / promo template (5–8 slides typical; this skeleton builds 4).

Duplicate ``build_examples_slide`` to add more use-case collage slides.

caw-slides variant: **宣伝・紹介・テストユーザー募集**

Use case
--------
Introducing a tool/program to an audience and recruiting users or testers.
Real screenshots are the primary evidence and the layout is a **denser
collage** than the research-talk variants. Follows the showcase rules in
``style-guide.md`` §15 (program-context header, short noun titles,
screenshot-led collage, caption-above-image, inline app-logo clusters).

Structure
---------
1. Concept slide  : program-context header + tool name; hero figure (left) +
                    two stacked explainer cards (right) + runtime app logos
2. Features slide : overview figure + card; two function diagrams below
3. Examples slide : screenshot collage, 2-line caption above each image,
                    inline app-logo clusters (実際に使ってきた証拠)
4. CTA slide      : recruitment ask

Replace every ``<...>`` placeholder. Drop real screenshots / logos into the
folder ``SHOT`` points at; until then the screenshot slots render as labelled
placeholder boxes and logo clusters are skipped (so the template runs as-is).

Usage
-----
1. Copy to ``office/presentation/scripts/generate_showcase_<YYYYMMDD>.py``.
   Also copy ``pptx_helpers.py`` and ``research_icons.py`` into the same folder
   (or keep the plugin installed so ``../references/`` resolves) — the locator
   below searches script-dir first, then ``../references/``.
2. Replace placeholders; set ``SHOT`` to your screenshot/logo folder
3. Run — output to ``presentations/slides/showcase_<YYYYMMDD>.pptx``
"""
from __future__ import annotations

import sys
from pathlib import Path

# Locate ``pptx_helpers.py``: script-dir first (after user copies it alongside),
# then ../references/ (running this script directly from the plugin install).
_HERE = Path(__file__).resolve().parent
for _cand in (_HERE, _HERE.parent / "references"):
    if (_cand / "pptx_helpers.py").is_file():
        _cand_str = str(_cand)
        if _cand_str not in sys.path:
            sys.path.insert(0, _cand_str)
        break

import pptx_helpers as h  # noqa: E402
import research_icons as ic  # noqa: E402
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

OUT_DIR = Path.cwd() / "presentations" / "slides"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_PATH = OUT_DIR / "showcase_template.pptx"
FIG_DIR = OUT_DIR / "_showcase_figs"
FIG_DIR.mkdir(parents=True, exist_ok=True)

# ─── Placeholders ─────────────────────────────────────────────────────────────
PROGRAM = "<プログラム/文脈ラベル>"            # 上部ヘッダ（例: 配布プログラム名）
TOOL_NAME = "ツール名：<Full Tool Name (abbr)>"  # 略称初出は full form で
SHOT = Path("<screenshot/logo folder>")        # 実スクショ・ロゴの置き場に変更する
RUNTIME_LOGOS: list[str] = []                  # 例: [str(SHOT/"a.png"), str(SHOT/"b.png")]


def _shot(slide, img, label, *, left, top, width, height):
    """Place a real screenshot if available, else a labelled placeholder box."""
    if img and Path(img).is_file():
        h.add_picture_fit(slide, img, left=left, top=top, max_width=width, max_height=height)
    else:
        h.add_shape_card(
            slide, left=left, top=top, width=width, height=height,
            paragraphs=[h.Paragraph(
                h.mixed_runs(label, size=Pt(16), color=h.COLOR_TEXT_BODY),
                alignment=PP_ALIGN.CENTER)],
            border=h.COLOR_TITLE, anchor=MSO_ANCHOR.MIDDLE,
        )
    return (left, top, width, height, f"<shot:{label}>")


# ─── Figures (research_icons; generic placeholders) ────────────────────────────

def _concept_fig():
    return ic.hub_diagram(
        "<ツール>",
        [("<部署1>", "<役割>", ic.icon_gear, ic.BLUE),
         ("<部署2>", "<役割>", ic.icon_flask, ic.GREEN),
         ("<部署3>", "<役割>", ic.icon_document, ic.ORANGE),
         ("<部署4>", "<役割>", ic.icon_slides, ic.AMBER)],
        FIG_DIR / "concept.png", center_sub="<サブ>")


def _features_fig():
    return ic.hub_diagram(
        "<4 機能>",
        [("<機能1>", None, ic.icon_gear, ic.BLUE),
         ("<機能2>", None, ic.icon_magnifier, ic.RED),
         ("<機能3>", None, ic.icon_molecule, ic.PURPLE),
         ("<機能4>", None, ic.icon_slides, ic.AMBER)],
        FIG_DIR / "features.png")


def _loop_fig():
    return ic.cycle_diagram(
        [("<段階1>", None, ic.icon_document, ic.ORANGE),
         ("<段階2>", None, ic.icon_gear, ic.BLUE),
         ("<段階3>", None, ic.icon_chart, ic.CYAN),
         ("<段階4>", None, ic.icon_magnifier, ic.PURPLE)],
        FIG_DIR / "loop.png", center_label="<改善>", center_sub="<サイクル>")


def _outputs_fig():
    return ic.hub_diagram(
        "<成果>",
        [("<形式1>", None, ic.icon_document, ic.ORANGE),
         ("<形式2>", None, ic.icon_document, ic.BLUE),
         ("<形式3>", None, ic.icon_slides, ic.AMBER)],
        FIG_DIR / "outputs.png")


def _cta_fig():
    return ic.converging_diagram(
        (ic.icon_gear, ic.NAVY),
        [("<参加者>", ic.icon_researcher, ic.BLUE),
         ("<参加者>", ic.icon_researcher, ic.GREEN),
         ("<参加者>", ic.icon_researcher, ic.ORANGE)],
        FIG_DIR / "cta.png", center_label="<ツール>")


# ─── Slides ────────────────────────────────────────────────────────────────────

def build_concept_slide(prs):
    slide = h.blank_slide(prs)
    rects = h.add_context_header(slide, PROGRAM, 1, tool_name=TOOL_NAME)

    hero = (Inches(0.7), Inches(1.55), Inches(5.9), Inches(4.55))
    h.add_picture_fit(slide, _concept_fig(), left=hero[0], top=hero[1],
                      max_width=hero[2], max_height=hero[3])
    rects.append(hero + ("<hero>",))

    c1 = (Inches(6.76), Inches(1.6), Inches(6.2), Inches(2.25))
    h.add_shape_card(slide, left=c1[0], top=c1[1], width=c1[2], height=c1[3],
                     paragraphs=[
                         h.Paragraph(h.mixed_runs("▸ <ツール> とは", size=Pt(24), bold=True, color=h.COLOR_TITLE)),
                         h.Paragraph(h.mixed_runs("・<一文で何をするツールか>", size=Pt(20))),
                         h.Paragraph(h.mixed_runs("・<誰のどんな課題を解くか>", size=Pt(20), bold=True, color=h.COLOR_EMPH_BLUE))],
                     anchor=MSO_ANCHOR.MIDDLE)
    rects.append(c1 + ("<card1>",))

    c2 = (Inches(6.76), Inches(4.0), Inches(6.2), Inches(2.0))
    h.add_shape_card(slide, left=c2[0], top=c2[1], width=c2[2], height=c2[3],
                     paragraphs=[
                         h.Paragraph(h.mixed_runs("▸ 仕組み", size=Pt(24), bold=True, color=h.COLOR_TITLE)),
                         h.Paragraph(h.mixed_runs("・<入口（窓口）の説明>", size=Pt(20))),
                         h.Paragraph(h.mixed_runs("・<拡張・育つ仕組み>", size=Pt(20)))],
                     anchor=MSO_ANCHOR.MIDDLE)
    rects.append(c2 + ("<card2>",))

    # 動作環境（基盤 CLI 等）+ ロゴ。RUNTIME_LOGOS が空ならロゴは描かれない。
    rt = (Inches(6.83), Inches(6.3), Inches(4.0), Inches(0.45))
    h.add_rich_text_box(slide,
                        [h.Paragraph(h.mixed_runs("動作環境: <CLI 名>", size=Pt(14), color=h.COLOR_TEXT_BODY))],
                        left=rt[0], top=rt[1], width=rt[2], height=rt[3])
    rects.append(rt + ("<runtime>",))
    rects.append(h.add_logo_cluster(slide, RUNTIME_LOGOS,
                                    left=Inches(10.95), top=Inches(6.1), width=Inches(2.0), height=Inches(0.75)))
    h.assert_no_overlap(rects)


def build_features_slide(prs, slide_number: int):
    slide = h.blank_slide(prs)
    rects = list(h.add_slide_chrome(slide, "メイン機能", slide_number))

    top_fig = (Inches(0.5), Inches(1.05), Inches(6.5), Inches(2.75))
    h.add_picture_fit(slide, _features_fig(), left=top_fig[0], top=top_fig[1],
                      max_width=top_fig[2], max_height=top_fig[3])
    rects.append(top_fig + ("<overview>",))

    card = (Inches(7.2), Inches(1.05), Inches(5.7), Inches(2.75))
    h.add_shape_card(slide, left=card[0], top=card[1], width=card[2], height=card[3],
                     paragraphs=[
                         h.Paragraph(h.mixed_runs("▸ 4 つの機能", size=Pt(24), bold=True, color=h.COLOR_TITLE)),
                         h.Paragraph(h.mixed_runs("・<機能1：一言>", size=Pt(20))),
                         h.Paragraph(h.mixed_runs("・<機能2：一言>", size=Pt(20))),
                         h.Paragraph(h.mixed_runs("・<機能3：一言>", size=Pt(20))),
                         h.Paragraph(h.mixed_runs("・<機能4：一言>", size=Pt(20)))],
                     anchor=MSO_ANCHOR.MIDDLE)
    rects.append(card + ("<card>",))

    rects.append(h.add_collage_caption(slide, "<機能A>", "<一言>",
                                       left=Inches(0.5), top=Inches(3.85), width=Inches(6.0), height=Inches(0.6)))
    fa = (Inches(0.5), Inches(4.5), Inches(6.0), Inches(2.7))
    h.add_picture_fit(slide, _loop_fig(), left=fa[0], top=fa[1], max_width=fa[2], max_height=fa[3])
    rects.append(fa + ("<figA>",))

    rects.append(h.add_collage_caption(slide, "<機能B>", "<一言>",
                                       left=Inches(6.9), top=Inches(3.85), width=Inches(6.0), height=Inches(0.6)))
    fb = (Inches(6.9), Inches(4.5), Inches(6.0), Inches(2.7))
    h.add_picture_fit(slide, _outputs_fig(), left=fb[0], top=fb[1], max_width=fb[2], max_height=fb[3])
    rects.append(fb + ("<figB>",))
    h.assert_no_overlap(rects)


def build_examples_slide(prs, slide_number: int):
    slide = h.blank_slide(prs)
    rects = list(h.add_slide_chrome(slide, "使用例", slide_number))

    cols = [(Inches(0.4), "<用途1>", "<一文説明>", h.COLOR_SUB_GREEN),
            (Inches(4.6), "<用途2>", "<一文説明>", h.COLOR_EMPH_BLUE),
            (Inches(8.8), "<用途3>", "<一文説明>", h.COLOR_TITLE)]
    col_w = Inches(4.0)
    for left, heading, sub, color in cols:
        rects.append(h.add_collage_caption(slide, heading, sub, left=left, top=Inches(1.4),
                                           width=col_w, color=color))
        rects.append(_shot(slide, None, heading, left=left, top=Inches(2.2),
                           width=col_w, height=Inches(3.1)))
        rects.append(h.add_logo_cluster(slide, [], left=left, top=Inches(5.45),
                                        width=col_w, height=Inches(0.5)))

    h.add_key_message_band(slide, [h.Paragraph(h.mixed_runs(
        "<L1: 実際にこれだけ幅広く使ってきた>", size=Pt(24), bold=True,
        color=h.COLOR_TEXT_BODY), alignment=PP_ALIGN.CENTER)])
    rects.append(h.CODEX_KEY_MSG_RECT + ("<key-msg>",))
    h.assert_no_overlap(rects)


def build_cta_slide(prs, slide_number: int):
    slide = h.blank_slide(prs)
    rects = list(h.add_slide_chrome(slide, "テストユーザー募集", slide_number))

    fig = (Inches(0.5), Inches(1.3), Inches(6.6), Inches(4.6))
    h.add_picture_fit(slide, _cta_fig(), left=fig[0], top=fig[1], max_width=fig[2], max_height=fig[3])
    rects.append(fig + ("<cta-fig>",))

    card = (Inches(7.3), Inches(1.4), Inches(5.6), Inches(4.4))
    h.add_shape_card(slide, left=card[0], top=card[1], width=card[2], height=card[3],
                     paragraphs=[
                         h.Paragraph(h.mixed_runs("▸ 募集要項", size=Pt(24), bold=True, color=h.COLOR_TITLE)),
                         h.Paragraph(h.mixed_runs("・対象：<どんな人に試してほしいか>", size=Pt(20))),
                         h.Paragraph(h.mixed_runs("・提供：<導入手順 / サポート>", size=Pt(20))),
                         h.Paragraph(h.mixed_runs("・始め方：<連絡方法>", size=Pt(20), bold=True, color=h.COLOR_EMPH_RED)),
                         h.Paragraph(h.mixed_runs("・<フィードバックのお願い>", size=Pt(20)))],
                     anchor=MSO_ANCHOR.MIDDLE)
    rects.append(card + ("<cta-card>",))

    h.add_key_message_band(slide, [h.Paragraph(h.mixed_runs(
        "<L1: まずは気軽に試してみませんか>", size=Pt(24), bold=True,
        color=h.COLOR_TEXT_BODY), alignment=PP_ALIGN.CENTER)])
    rects.append(h.CODEX_KEY_MSG_RECT + ("<key-msg>",))
    h.assert_no_overlap(rects)


def main() -> None:
    prs = h.new_presentation()
    build_concept_slide(prs)
    build_features_slide(prs, slide_number=2)
    build_examples_slide(prs, slide_number=3)
    build_cta_slide(prs, slide_number=4)
    prs.save(str(OUT_PATH))
    print(f"Saved {OUT_PATH}")


if __name__ == "__main__":
    main()
