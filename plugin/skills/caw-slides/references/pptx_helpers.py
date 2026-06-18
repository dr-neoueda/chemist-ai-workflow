"""Shared helpers for generating research presentation slides.

Encodes the caw-slides style guide (see ``references/style-guide.md``) so individual
``generate_<purpose>_<YYYYMMDD>.py`` scripts stay short and consistent.

All public functions return the created object so callers can further
tweak it when the style guide does not cover an edge case.

Layout: 16:9 (13.33" x 7.5"). Bilingual font split: 和文 = Noto Sans JP (fallback: MS Gothic),
英数字 = Segoe UI (fallback: Arial; auto-detected per character by ``mixed_runs`` / ``Run.font=None``).
Drop-in to any chemistry project — no project-specific names or paths.
"""
from __future__ import annotations

import os
import tempfile
import warnings
from dataclasses import dataclass
from pathlib import Path

import matplotlib.font_manager as fm
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Fonts
# ---------------------------------------------------------------------------

# Preferred fonts (参考デザイン: 日本語 = Noto Sans JP / 英数字 = Segoe UI).
# python-pptx merely *names* these in the .pptx; PowerPoint substitutes the
# installed face if the preferred one is absent. The *_FALLBACK names document
# the guaranteed-available faces (Windows / bundled with MS Office) that should
# stand in. matplotlib resolves its own font *file* via ``JP_FONT_PATH`` below
# (Noto Sans JP → MS Gothic), independent of these name constants.
FONT_JA = "Noto Sans JP"
FONT_EN = "Segoe UI"
FONT_JA_FALLBACK = "MS Gothic"
FONT_EN_FALLBACK = "Arial"

# Cross-platform MS Gothic candidates. First-existing path wins.
# Override with the ``CAW_SLIDES_MSGOTHIC`` environment variable if needed.
_MSGOTHIC_CANDIDATES: tuple[Path, ...] = (
    # macOS — bundled with Microsoft PowerPoint
    Path("/Applications/Microsoft PowerPoint.app/Contents/Resources/DFonts/msgothic.ttc"),
    # macOS — bundled with Microsoft Word (also valid)
    Path("/Applications/Microsoft Word.app/Contents/Resources/DFonts/msgothic.ttc"),
    # Windows
    Path("C:/Windows/Fonts/msgothic.ttc"),
    Path("C:/Windows/Fonts/MSGOTHIC.TTC"),
    # Linux — common install paths if user has set up MS Gothic manually
    Path("/usr/share/fonts/truetype/msgothic/msgothic.ttc"),
    Path("/usr/share/fonts/MS/msgothic.ttc"),
)


def _find_msgothic_path() -> Path | None:
    """Locate MS Gothic on disk. Returns ``None`` if it is unavailable.

    Honors the ``CAW_SLIDES_MSGOTHIC`` env var first so a user can point to a
    non-standard install (``export CAW_SLIDES_MSGOTHIC=/path/to/msgothic.ttc``).
    """
    env_override = os.environ.get("CAW_SLIDES_MSGOTHIC")
    if env_override:
        candidate = Path(env_override)
        if candidate.is_file():
            return candidate
    for candidate in _MSGOTHIC_CANDIDATES:
        if candidate.is_file():
            return candidate
    return None


# Resolved at import time. ``None`` on platforms without MS Gothic — callers
# that need Japanese in matplotlib output should detect this and either
# install MS Gothic or fall back to a CJK-capable system font.
MSGOTHIC_PATH = _find_msgothic_path()

# Noto Sans JP candidates (preferred 日本語 face). Override with the
# ``CAW_SLIDES_JPFONT`` environment variable.
_NOTO_JP_CANDIDATES: tuple[Path, ...] = (
    Path.home() / "Library/Fonts/NotoSansJP-Regular.otf",
    Path.home() / "Library/Fonts/NotoSansJP-Regular.ttf",
    Path("/Library/Fonts/NotoSansJP-Regular.otf"),
    Path("/Library/Fonts/NotoSansCJKjp-Regular.otf"),
    Path("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"),
    Path("/usr/share/fonts/truetype/noto/NotoSansJP-Regular.otf"),
    Path("/usr/share/fonts/google-noto-cjk/NotoSansCJK-Regular.ttc"),
    Path("C:/Windows/Fonts/NotoSansJP-Regular.otf"),
)


def _find_jp_font_path() -> Path | None:
    """Locate a Japanese-capable font file, preferring Noto Sans JP, then MS Gothic.

    Honors ``CAW_SLIDES_JPFONT`` then ``CAW_SLIDES_MSGOTHIC`` env overrides so a
    user can point at a non-standard install. Returns ``None`` if neither face is
    available (callers fall back to ASCII-only output / raise).
    """
    for env_var in ("CAW_SLIDES_JPFONT", "CAW_SLIDES_MSGOTHIC"):
        override = os.environ.get(env_var)
        if override and Path(override).is_file():
            return Path(override)
    for candidate in (*_NOTO_JP_CANDIDATES, *_MSGOTHIC_CANDIDATES):
        if candidate.is_file():
            return candidate
    return None


# Preferred JP font file for matplotlib (Noto Sans JP → MS Gothic fallback).
JP_FONT_PATH = _find_jp_font_path()


def configure_matplotlib_japanese() -> str:
    """Register the preferred JP font (Noto Sans JP → MS Gothic) with matplotlib.

    Also sets the §0「全て太字」defaults: bold tick labels, axis labels and
    titles so figure text reads on a projector.

    Returns
    -------
    str
        The matplotlib font-family name that was activated.

    Raises
    ------
    FileNotFoundError
        If neither Noto Sans JP nor MS Gothic is found on any default path and
        neither ``CAW_SLIDES_JPFONT`` nor ``CAW_SLIDES_MSGOTHIC`` is set.

    Notes
    -----
    Must be called before any matplotlib rendering that includes Japanese
    text. Without this, characters render as tofu (縦長の □).
    """
    if JP_FONT_PATH is None:
        raise FileNotFoundError(
            "No Japanese font found. Install Noto Sans JP, or Microsoft "
            "PowerPoint (bundles MS Gothic), or set CAW_SLIDES_JPFONT / "
            "CAW_SLIDES_MSGOTHIC to a .ttc/.otf/.ttf path."
        )
    fm.fontManager.addfont(str(JP_FONT_PATH))
    family = fm.FontProperties(fname=str(JP_FONT_PATH)).get_name()
    plt.rcParams["font.family"] = family
    # §0「全て太字」: figure text (labels, ticks, titles) bold for projector legibility
    plt.rcParams["font.weight"] = "bold"
    plt.rcParams["axes.labelweight"] = "bold"
    plt.rcParams["axes.titleweight"] = "bold"
    plt.rcParams["axes.spines.top"] = False
    plt.rcParams["axes.spines.right"] = False
    return family


# ---------------------------------------------------------------------------
# Palette (see CLAUDE.md §3)
# ---------------------------------------------------------------------------

# --- Base color (参考デザイン: teal #3686A6 = RGB 54,134,166) + モノクロ階調 ---
COLOR_BASE = RGBColor(0x36, 0x86, 0xA6)            # base teal
COLOR_BASE_DARK = RGBColor(0x1E, 0x4E, 0x63)       # 暗: タイトルバー
COLOR_BASE_DARKER = RGBColor(0x14, 0x35, 0x44)     # 最暗
COLOR_BASE_LIGHT = RGBColor(0xA9, 0xD0, 0xE0)      # 明: 小見出しバー
COLOR_BASE_LIGHTER = RGBColor(0xDC, 0xEC, 0xF2)    # 最明: key-message 箱
# --- Accents (参考デザイン) ---
COLOR_ACCENT_TERRACOTTA = RGBColor(0xBA, 0x59, 0x36)  # 186,89,54  暖色強調・データ
COLOR_ACCENT_GREEN = RGBColor(0x45, 0x9B, 0x2D)       # 69,155,45  第2強調
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

COLOR_TEXT_BODY = RGBColor(0x22, 0x22, 0x22)
# Role colors — 参考デザイン体系に再編（名前は据え置き＝既存テンプレ互換）
COLOR_TITLE = COLOR_BASE_DARK                    # 見出し・構造（白地のタイトル文字/▸見出し）
COLOR_EMPH_BLUE = COLOR_BASE                     # 要点（base teal）
COLOR_EMPH_NAVY = COLOR_BASE_DARK
COLOR_EMPH_RED = COLOR_ACCENT_TERRACOTTA         # 注意・暖色強調（純赤 → テラコッタ）
COLOR_SUB_CYAN = COLOR_BASE_LIGHT
COLOR_SUB_GREEN = COLOR_ACCENT_GREEN             # 第2強調（緑 #459B2D）

COLOR_ACCENT_BLUE = COLOR_BASE                   # 構造色（テーブルヘッダ・バー等）
COLOR_ACCENT_LIGHTBLUE = COLOR_BASE_LIGHT
COLOR_ACCENT_CYAN = COLOR_BASE_LIGHT
COLOR_ACCENT_RED = COLOR_ACCENT_TERRACOTTA
COLOR_HIGHLIGHT_YELLOW = RGBColor(0xFF, 0xFF, 0x00)
COLOR_HIGHLIGHT_AMBER = RGBColor(0xFF, 0xC0, 0x00)

# 7原則の強調色 (§12)
COLOR_EMPH_7PRINCIPLES_RED = COLOR_ACCENT_TERRACOTTA

# ---------------------------------------------------------------------------
# Categorical palette — 図表・グラフ・多系列で青一色を避け、識別しやすい配色に
# (§3 色分けルール準拠: 対比は赤/青/緑、系列はこの順で巡回)
# ---------------------------------------------------------------------------
# matplotlib 用（hex 文字列）
CATEGORICAL_HEX: tuple[str, ...] = (
    "#3686A6",  # base teal（参考デザイン）
    "#BA5936",  # テラコッタ
    "#459B2D",  # 緑
    "#1E4E63",  # 暗 teal
    "#D9A441",  # アンバー
    "#7A5BA6",  # 紫
    "#A0A0A0",  # グレー
)
# native chart 用（RGBColor）
CATEGORICAL_RGB: tuple[RGBColor, ...] = tuple(
    RGBColor(int(hx[1:3], 16), int(hx[3:5], 16), int(hx[5:7], 16)) for hx in CATEGORICAL_HEX
)
# 行ハイライト用の淡色塗り（テーブルの hero 行など）
COLOR_ROW_HIGHLIGHT_FILL = RGBColor(0xFF, 0xF2, 0xCC)  # 淡いアンバー

# ---------------------------------------------------------------------------
# Layout constants (16:9 slide, §1)
# ---------------------------------------------------------------------------

SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

MARGIN_LEFT = Inches(0.4)
TITLE_TOP = Inches(0.12)
TITLE_HEIGHT = Inches(0.7)
SEPARATOR_Y = Inches(0.82)
BODY_TOP = Inches(0.95)
BODY_HEIGHT = Inches(5.6)
KEY_MESSAGE_Y = Inches(6.8)

# Font sizes (§2, §12)
SIZE_TITLE_MAIN = Pt(72)
SIZE_THEME = Pt(32)
SIZE_SLIDE_TITLE = Pt(28)
SIZE_SECTION = Pt(24)
SIZE_BODY = Pt(20)
SIZE_EMPH = Pt(24)
SIZE_NOTE = Pt(12)


# ---------------------------------------------------------------------------
# Presentation bootstrap
# ---------------------------------------------------------------------------


def new_presentation() -> Presentation:
    """Create a blank 16:9 presentation sized per the style guide."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def blank_slide(prs: Presentation):
    """Add a blank slide (layout 6 is the blank layout in default master)."""
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------------------------------------------------------------------------
# Text helpers
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class Run:
    """A single styled text run. Use multiple Run objects to mix JA/EN."""

    text: str
    size: Pt = SIZE_BODY
    bold: bool = False
    color: RGBColor = COLOR_TEXT_BODY
    # None = auto: JA if any Hiragana/Katakana/Kanji, else EN
    font: str | None = None


def _is_japanese(text: str) -> bool:
    for ch in text:
        code = ord(ch)
        if (
            0x3000 <= code <= 0x303F  # CJK symbols/punctuation (、 。 「 」 …)
            or 0x3040 <= code <= 0x309F  # Hiragana
            or 0x30A0 <= code <= 0x30FF  # Katakana
            or 0x4E00 <= code <= 0x9FFF  # CJK Unified Ideographs
            or 0xFF00 <= code <= 0xFFEF  # Halfwidth / Fullwidth forms
        ):
            return True
    return False


def add_text_box(
    slide,
    runs: list[Run],
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    word_wrap: bool = True,
    auto_size: bool = True,
):
    """Add a textbox whose runs are individually font-styled.

    Parameters
    ----------
    runs
        Sequence of :class:`Run`. Mix JA/EN runs to satisfy §1 (MS Gothic
        for 和文, Arial for 英数字) within a single textbox.
    auto_size
        If True, enables text-autofit-to-shape. **Do not rely on this to
        fit arbitrary amounts of text** — see §12 layout verification.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = anchor
    if auto_size:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    for idx, run_spec in enumerate(runs):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = run_spec.text
        font = run.font
        font.size = run_spec.size
        font.bold = run_spec.bold
        font.color.rgb = run_spec.color
        font.name = run_spec.font or (
            FONT_JA if _is_japanese(run_spec.text) else FONT_EN
        )

    return box


def add_title(slide, text: str, *, slide_number: int | None = None):
    """Add the standard slide title + blue separator line (§12 layout)."""
    add_text_box(
        slide,
        [Run(text, size=SIZE_SLIDE_TITLE, bold=True, color=COLOR_TITLE)],
        left=MARGIN_LEFT,
        top=TITLE_TOP,
        width=Inches(11.0),
        height=TITLE_HEIGHT,
        auto_size=False,
    )
    line = slide.shapes.add_connector(
        1, MARGIN_LEFT, SEPARATOR_Y, SLIDE_WIDTH - MARGIN_LEFT, SEPARATOR_Y
    )
    line.line.color.rgb = COLOR_ACCENT_BLUE
    line.line.width = Pt(1.5)

    if slide_number is not None:
        add_text_box(
            slide,
            [Run(str(slide_number), size=SIZE_NOTE, color=COLOR_TEXT_BODY)],
            left=SLIDE_WIDTH - Inches(0.8),
            top=TITLE_TOP,
            width=Inches(0.4),
            height=Inches(0.4),
            align=PP_ALIGN.RIGHT,
        )


def add_key_message(slide, runs: list[Run]):
    """Add the bottom key-message band (§12 layout)."""
    add_text_box(
        slide,
        runs,
        left=MARGIN_LEFT,
        top=KEY_MESSAGE_Y,
        width=SLIDE_WIDTH - 2 * MARGIN_LEFT,
        height=Inches(0.55),
        auto_size=False,
    )


# ---------------------------------------------------------------------------
# Image helpers
# ---------------------------------------------------------------------------


def add_picture_fit(
    slide,
    image_path: str | Path,
    *,
    left: Emu,
    top: Emu,
    max_width: Emu,
    max_height: Emu,
):
    """Add a picture preserving aspect ratio, centered in a bounding box.

    The **final** rendered width/height — not ``max_width``/``max_height``
    — is what you should reason about when checking for overlap with
    neighboring shapes (§12 layout verification).
    """
    from PIL import Image

    with Image.open(image_path) as img:
        iw, ih = img.size
    # Convert pixel dimensions to EMU assuming 96 DPI (1 px = 9525 EMU)
    px_w = Emu(int(iw * 9525))
    px_h = Emu(int(ih * 9525))
    ratio = min(max_width / px_w, max_height / px_h)
    target_w = Emu(int(iw * 9525 * ratio))
    target_h = Emu(int(ih * 9525 * ratio))
    x = left + (max_width - target_w) // 2
    y = top + (max_height - target_h) // 2
    return slide.shapes.add_picture(
        str(image_path), x, y, width=target_w, height=target_h
    )


# ---------------------------------------------------------------------------
# Showcase / 宣伝・紹介デッキ helpers (§15)
# ---------------------------------------------------------------------------


def add_context_header(
    slide,
    program_label: str,
    slide_number: int,
    *,
    tool_name: str | None = None,
) -> list[Rect]:
    """showcase（宣伝・紹介・募集）デッキ用ヘッダ。

    研究発表用の :func:`add_slide_chrome` がスライド固有の断定見出しを置くのに対し、
    こちらはプログラム/文脈ラベル（例: 配布プログラム名）をタイトル位置に置く想定。
    ``tool_name`` を渡すと区切り線の下にツール正式名称を中央サブ行で添える
    （タイトルスライド向け。略称初出は full form で書くこと）。

    Returns rects (chrome + optional tool-name) for :func:`assert_no_overlap`.
    """
    rects = list(add_slide_chrome(slide, program_label, slide_number))
    if tool_name is not None:
        tn = (Inches(2.5), Inches(0.92), Inches(8.33), Inches(0.5))
        add_rich_text_box(
            slide,
            [Paragraph(mixed_runs(tool_name, size=Pt(18), color=COLOR_TEXT_BODY),
                       alignment=PP_ALIGN.CENTER)],
            left=tn[0], top=tn[1], width=tn[2], height=tn[3],
            anchor=MSO_ANCHOR.MIDDLE,
        )
        rects.append(tn + ("<tool-name>",))
    return rects


def add_collage_caption(
    slide,
    heading: str,
    sub: str,
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    color: RGBColor = COLOR_TITLE,
    height: Emu = Inches(0.67),
) -> Rect:
    """コラージュ用 2 行キャプション（スクリーンショットの真上に置く）。

    1 行目 = 見出し（``color`` で着色・bold）、2 行目 = 一文説明（本文色）。
    showcase デッキの使用例スライドで、各画像の上に左寄せで置く想定。
    Returns the caption rect for :func:`assert_no_overlap`.
    """
    add_rich_text_box(
        slide,
        [Paragraph(mixed_runs(heading, size=Pt(20), bold=True, color=color)),
         Paragraph(mixed_runs(sub, size=Pt(14), color=COLOR_TEXT_BODY))],
        left=left, top=top, width=width, height=height,
    )
    return (left, top, width, height, f"<caption:{heading}>")


def add_logo_cluster(
    slide,
    icon_paths: list[str | Path],
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu = Inches(0.5),
    slot: Emu = Inches(0.6),
) -> Rect:
    """使用アプリのロゴを小さく等間隔で 1 列に並べる（実使用ツールの裏付け）。

    showcase デッキで使用例キャプションの脇・画像下に置く想定。各ロゴはアスペクト比を
    保って ``slot`` 幅のスロット内に中央寄せされる。``icon_paths`` が空、または割り当て幅が
    0 以下になる場合は何も描かない。

    返すのは **宣言した** bounding rect ``(left, top, width, height)`` で、実際に描画される
    ロゴ群の footprint より広い（中央寄せの余白を含む）。これは :func:`add_picture_fit` と
    同じ方針で、:func:`assert_no_overlap` には保守的（広め）に効く。
    """
    n = len(icon_paths)
    if n == 0:
        return (left, top, width, height, "<logos:empty>")
    slot = min(slot, width // n)
    if slot <= 0:
        return (left, top, width, height, "<logos:empty>")
    x0 = left + (width - slot * n) // 2
    for i, path in enumerate(icon_paths):
        add_picture_fit(
            slide, path, left=x0 + slot * i, top=top,
            max_width=slot, max_height=height,
        )
    return (left, top, width, height, "<logos>")


# ---------------------------------------------------------------------------
# Forbidden-character guard (§11)
# ---------------------------------------------------------------------------

FORBIDDEN_IN_MATPLOTLIB = {
    "✓": "○",  # ✓
    "✗": "×",  # ✗
    "−": "-",  # Unicode minus
}
# Unicode subscript digits U+2080..U+2089 → ASCII
for _i in range(10):
    FORBIDDEN_IN_MATPLOTLIB[chr(0x2080 + _i)] = str(_i)


def sanitize_for_matplotlib(text: str) -> str:
    """Replace glyphs missing from MS Gothic with safe ASCII fallbacks."""
    for bad, good in FORBIDDEN_IN_MATPLOTLIB.items():
        text = text.replace(bad, good)
    return text


# ---------------------------------------------------------------------------
# Layout overlap guard (CLAUDE.md §「テキストボックス重なり禁止」)
# ---------------------------------------------------------------------------

Rect = tuple[int, int, int, int, str]


def assert_no_overlap(rects: list[Rect]) -> None:
    """Raise ``ValueError`` if any two rectangles in ``rects`` overlap.

    Each rectangle is ``(left, top, width, height, label)`` in EMU. Tangent
    edges (one rect ending exactly where the next begins) are **not** counted
    as overlapping -- honest adjacency is fine, visible intersection is not.
    """
    for i, a in enumerate(rects):
        ax, ay, aw, ah, al = a
        for b in rects[i + 1 :]:
            bx, by, bw, bh, bl = b
            if ax + aw <= bx or bx + bw <= ax:
                continue
            if ay + ah <= by or by + bh <= ay:
                continue
            raise ValueError(
                f"Layout overlap: {al!r} at ({ax},{ay},{aw}x{ah}) "
                f"overlaps {bl!r} at ({bx},{by},{bw}x{bh})"
            )


# ---------------------------------------------------------------------------
# Table helper (prefer tables over bullet lists for comparison data)
# ---------------------------------------------------------------------------


# ---------------------------------------------------------------------------
# Codex-style layout (2026-04-24: adopted from reference decks)
# ---------------------------------------------------------------------------

COLOR_CARD_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_CARD_BORDER = RGBColor(0xDD, 0xDD, 0xDD)
COLOR_PILL_BLUE_FILL = RGBColor(0xDC, 0xEC, 0xF2)   # 淡 teal（参考デザイン）
COLOR_PILL_RED_FILL = RGBColor(0xF4, 0xE3, 0xDC)    # 淡テラコッタ
COLOR_PILL_GREEN_FILL = RGBColor(0xE4, 0xF1, 0xDD)  # 淡グリーン
COLOR_PILL_GREY_FILL = RGBColor(0xF4, 0xF6, 0xF8)
COLOR_KEY_MSG_FILL = RGBColor(0xDC, 0xEC, 0xF2)     # key-message 箱 = 最明 teal
# Source-line colour is now quite light and the font is small — the source
# line must read as marginalia, not content.
COLOR_SOURCE_GREY = RGBColor(0xB0, 0xB0, 0xB0)
COLOR_LABEL_GREY = RGBColor(0x88, 0x88, 0x88)   # small caption labels (≤ 12pt)

SIZE_CARD_SECTION = Pt(21)  # ▸ headings inside body card
SIZE_CARD_BULLET = Pt(16)   # ・ bullet text
SIZE_PILL_TITLE = Pt(14)
SIZE_PILL_NUMBER = Pt(19)
SIZE_PILL_CAPTION = Pt(12)

# Fixed reference positions (copied from Codex reference decks)
CODEX_TITLE_RECT = (Inches(0.4), Inches(0.12), Inches(11.7), Inches(0.45))
CODEX_SLIDE_NUM_RECT = (Inches(12.18), Inches(0.12), Inches(0.62), Inches(0.4))
CODEX_SEP_RECT = (Inches(0.4), Inches(0.82), Inches(12.53), Inches(0.03))  # deprecated: 新 add_slide_chrome は塗りバーで separator 線を描かない（互換用に残置・overlap rects に足さない）
CODEX_KEY_MSG_RECT = (Inches(0.52), Inches(6.28), Inches(12.25), Inches(0.58))
CODEX_SOURCE_RECT = (Inches(0.52), Inches(7.02), Inches(12.2), Inches(0.35))


@dataclass(frozen=True)
class Paragraph:
    """A single text paragraph composed of multiple styled runs.

    Use this when a paragraph needs to mix Japanese (MS Gothic) and English
    (Arial) runs, or to emphasize inline keywords in a different color.
    """

    runs: list["Run"]
    alignment: PP_ALIGN = PP_ALIGN.LEFT


def _apply_runs_to_paragraph(paragraph, runs: list["Run"]) -> None:
    for run_spec in runs:
        r = paragraph.add_run()
        r.text = run_spec.text
        r.font.size = run_spec.size
        r.font.bold = run_spec.bold
        r.font.color.rgb = run_spec.color
        r.font.name = run_spec.font or (
            FONT_JA if _is_japanese(run_spec.text) else FONT_EN
        )


def _write_paragraphs(text_frame, paragraphs: list[Paragraph]) -> None:
    for p_idx, para in enumerate(paragraphs):
        if p_idx == 0:
            p = text_frame.paragraphs[0]
            # Clear any default run
            for run in list(p.runs):
                run.text = ""
        else:
            p = text_frame.add_paragraph()
        p.alignment = para.alignment
        _apply_runs_to_paragraph(p, para.runs)


def add_rich_text_box(
    slide,
    paragraphs: list[Paragraph],
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    word_wrap: bool = True,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    auto_size: bool = False,
):
    """Add a text box supporting multiple paragraphs, each with multiple runs.

    Use this in place of :func:`add_text_box` whenever a paragraph mixes
    JA/EN or has inline color/weight changes.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = anchor
    if auto_size:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    _write_paragraphs(tf, paragraphs)
    return box


def _style_shape_fill(shape, fill: RGBColor | None, border: RGBColor | None) -> None:
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if border is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border
        shape.line.width = Pt(0.75)


def _set_shape_shadow(shape, on: bool = False) -> None:
    """Render the shape **flat** — drop-shadow は使わない（スタイルガイド §12）.

    影はスライド全体で不使用の方針。PowerPoint テーマは矩形に既定で淡い
    drop-shadow を付けるため、空の ``<a:effectLst/>`` を注入してテーマ既定の
    影を必ず無効化する。``on`` 引数は呼び出し後方互換のために残すが**無視**する
    （常にフラット）。
    """
    from lxml import etree
    from pptx.oxml.ns import qn

    spPr = shape._element.spPr
    existing = spPr.find(qn("a:effectLst"))
    if existing is not None:
        spPr.remove(existing)
    etree.SubElement(spPr, qn("a:effectLst"))  # 空 = エフェクトなし（フラット）


def mixed_runs(
    text: str,
    *,
    size: Pt,
    bold: bool = False,
    color: RGBColor = COLOR_TEXT_BODY,
) -> list["Run"]:
    """Split a mixed-language string into per-language runs.

    Each character is classified as Japanese (MS Gothic) or English
    (Arial, which covers ASCII, digits, Latin-extended). Consecutive
    characters of the same class are merged into one Run so paragraph
    formatting stays clean.
    """
    if not text:
        return []
    parts: list[tuple[str, bool]] = []
    cur_text = text[0]
    cur_ja = _is_japanese(text[0])
    for ch in text[1:]:
        is_ja = _is_japanese(ch)
        if is_ja == cur_ja:
            cur_text += ch
        else:
            parts.append((cur_text, cur_ja))
            cur_text = ch
            cur_ja = is_ja
    parts.append((cur_text, cur_ja))
    return [
        Run(segment, size=size, bold=bold, color=color,
            font=FONT_JA if is_ja else FONT_EN)
        for segment, is_ja in parts
    ]


def add_shape_card(
    slide,
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    paragraphs: list[Paragraph] | None = None,
    fill: RGBColor | None = None,
    border: RGBColor | None = None,
    padding: Emu = Inches(0.2),
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    shadow: bool = False,
):
    """Add a rectangular card (AUTO_SHAPE) optionally filled with paragraphs.

    The card is a single shape — its text lives inside the card's text
    frame. Unlike the old ``add_text_box`` + decorative background split,
    this keeps text and panel as one shape so ``assert_no_overlap`` only
    sees one rectangle per card.

    ``padding`` is applied as internal margins so bullet text does not kiss
    the rounded edge of the card.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    _style_shape_fill(shape, fill, border)
    _set_shape_shadow(shape, shadow)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = padding
    tf.margin_right = padding
    tf.margin_top = padding
    tf.margin_bottom = padding
    if paragraphs:
        _write_paragraphs(tf, paragraphs)
    return shape


def add_pill(
    slide,
    paragraphs: list[Paragraph],
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    fill: RGBColor = COLOR_PILL_BLUE_FILL,
    border: RGBColor | None = None,
    padding: Emu = Inches(0.06),
    shadow: bool = False,
):
    """Small colored callout pill (AUTO_SHAPE). Text is centered by default.

    Always rendered **flat** (no drop-shadow): 影はスライド全体で不使用のため、
    ``shadow`` 引数は無視される（後方互換のためにのみ残置）。強調は塗り・枠・
    フォントサイズ／色で付ける。
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    _style_shape_fill(shape, fill, border)
    _set_shape_shadow(shape, shadow)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = padding
    tf.margin_right = padding
    tf.margin_top = padding
    tf.margin_bottom = padding
    _write_paragraphs(tf, paragraphs)
    return shape


def add_key_message_band(
    slide,
    paragraphs: list[Paragraph],
    *,
    fill: RGBColor = COLOR_KEY_MSG_FILL,
    border: RGBColor | None = None,
    shadow: bool = False,
):
    """Full-width key-message band at the Codex fixed position (y=6.28).

    Drawn as a **rounded rectangle with no outline and no shadow** (flat):
    淡ティールの塗りと角丸だけで主張を示す。影はスライド全体で不使用のため
    ``shadow`` 引数は無視される（後方互換のためにのみ残置）。
    """
    left, top, width, height = CODEX_KEY_MSG_RECT
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    _style_shape_fill(shape, fill, border)
    _set_shape_shadow(shape, shadow)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    _write_paragraphs(tf, paragraphs)
    return shape


def add_source_line(slide, text: str):
    """Tiny source credit at the Codex fixed position (y=7.02).

    Pt 9 + light grey (#B0B0B0) — marginalia only, must not compete with
    the slide's content for the reader's attention.
    """
    left, top, width, height = CODEX_SOURCE_RECT
    runs = [Run(text, size=Pt(9), bold=False, color=COLOR_SOURCE_GREY,
                font=FONT_EN)]
    return add_rich_text_box(
        slide, [Paragraph(runs)], left=left, top=top, width=width,
        height=height, anchor=MSO_ANCHOR.TOP,
    )


def add_category_pill(
    slide,
    text: str,
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    fill: RGBColor = COLOR_PILL_BLUE_FILL,
    border: RGBColor | None = None,
    shadow: bool = False,
):
    """Title-slide category pill (e.g. '報告会 No.12')."""
    runs = mixed_runs(text, size=Pt(20), bold=True, color=COLOR_TITLE)
    return add_pill(
        slide,
        [Paragraph(runs, alignment=PP_ALIGN.CENTER)],
        left=left, top=top, width=width, height=height,
        fill=fill, border=border, shadow=shadow,
    )


# ---------------------------------------------------------------------------
# Native flow diagram (editable PowerPoint shapes, slides 3 / 9)
# ---------------------------------------------------------------------------


def add_flow_box(
    slide,
    text: str,
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    fill: RGBColor,
    text_color: RGBColor = COLOR_TEXT_BODY,
    font_size: Pt = Pt(12),
    bold: bool = True,
    shadow: bool = False,
    shape_type: MSO_SHAPE = MSO_SHAPE.ROUNDED_RECTANGLE,
):
    """Single box for a flow/timeline diagram with centered multi-line text.

    ``text`` may contain ``\\n`` for line breaks; each line becomes its own
    paragraph so PowerPoint wraps cleanly. Language-aware font switching is
    applied per paragraph.
    """
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    _style_shape_fill(shape, fill, None)
    _set_shape_shadow(shape, shadow)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    paragraphs = [
        Paragraph(
            mixed_runs(line, size=font_size, bold=bold, color=text_color),
            alignment=PP_ALIGN.CENTER,
        )
        for line in text.split("\n")
    ]
    _write_paragraphs(tf, paragraphs)
    return shape


def add_flow_arrow(
    slide,
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    direction: str = "right",
    fill: RGBColor = COLOR_TEXT_BODY,
    border: RGBColor | None = None,
):
    """Add a filled arrow shape (``MSO_SHAPE.<direction>_ARROW``).

    A dedicated arrow auto-shape is used instead of a ``MSO_CONNECTOR``
    because connectors with zero-width/zero-height bounding boxes (purely
    horizontal or vertical) occasionally serialize coordinates as floats
    that break python-pptx round-trips. A filled auto-shape has an
    explicit rectangular bounding box and is trivially selectable /
    editable in PowerPoint.

    ``direction`` is one of ``"right"``, ``"left"``, ``"up"``, ``"down"``.
    """
    shape_map = {
        "right": MSO_SHAPE.RIGHT_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
    }
    shape_type = shape_map[direction]
    shape = slide.shapes.add_shape(
        shape_type,
        int(left), int(top), int(width), int(height),
    )
    _style_shape_fill(shape, fill, border)
    _set_shape_shadow(shape, False)
    # Arrows are decorative; do not add text.
    return shape


# ---------------------------------------------------------------------------
# Native Excel-editable charts (replaces matplotlib PNGs for graphs)
# ---------------------------------------------------------------------------


def _style_chart_common(
    chart,
    *,
    title: str | None,
    show_legend: bool = True,
    legend_font: Pt = Pt(11),
    title_font: Pt = Pt(13),
    axis_font: Pt = Pt(10),
    plot_area_border: RGBColor | None = RGBColor(0x00, 0x00, 0x00),
    plot_area_border_width: Pt = Pt(1.0),
    gridlines: bool = False,
    gridlines_color: RGBColor = RGBColor(0xDD, 0xDD, 0xDD),
    gridlines_width: Pt = Pt(0.5),
) -> None:
    """Apply consistent styling across chart types.

    Default behavior (per 2026-04-24 feedback):
    - **Plot-area border**: black 1pt (clear publication-style frame).
    - **Gridlines**: OFF. Data readability takes priority over reference
      lines. Pass ``gridlines=True`` to enable subtle thin pale-grey
      reference lines when a specific chart needs them.
    """
    if title:
        chart.has_title = True
        tframe = chart.chart_title.text_frame
        tframe.text = title
        for p in tframe.paragraphs:
            for r in p.runs:
                r.font.size = title_font
                r.font.bold = True
                r.font.name = FONT_JA if _is_japanese(r.text) else FONT_EN
    else:
        chart.has_title = False
    chart.has_legend = show_legend
    if show_legend:
        try:
            chart.legend.include_in_layout = False
            legend_font_obj = chart.legend.font
            legend_font_obj.size = legend_font
            legend_font_obj.name = FONT_EN
        except Exception:
            pass
    try:
        for axis in (chart.category_axis, chart.value_axis):
            if axis is None:
                continue
            axis.tick_labels.font.size = axis_font
            axis.tick_labels.font.name = FONT_EN
    except Exception:
        pass

    # Plot-area border: python-pptx 1.0.2 doesn't expose Chart.plot_area,
    # so we inject <c:spPr>/<a:ln> into <c:plotArea> via XML. The spPr must
    # be the last child of plotArea (just before an optional c:extLst) per
    # the OOXML schema.
    if plot_area_border is not None:
        try:
            from lxml import etree
            from pptx.oxml.ns import qn
            chart_space = chart.element  # <c:chartSpace>
            chart_xml = chart_space.find(qn("c:chart"))
            plot_area_xml = chart_xml.find(qn("c:plotArea"))
            # Remove any existing spPr so we replace cleanly
            for existing in plot_area_xml.findall(qn("c:spPr")):
                plot_area_xml.remove(existing)
            # Detach extLst if present, append new spPr, then re-append extLst
            extLst = plot_area_xml.find(qn("c:extLst"))
            if extLst is not None:
                plot_area_xml.remove(extLst)
            spPr = etree.SubElement(plot_area_xml, qn("c:spPr"))
            # Transparent plot-area fill so the data region reads as "paper"
            etree.SubElement(spPr, qn("a:noFill"))
            # a:ln w="..." expects EMU. pptx.util.Pt is already an int in
            # EMU (1 pt = 12700 EMU), so int(Pt(1.0)) == 12700 is correct.
            ln = etree.SubElement(spPr, qn("a:ln"))
            ln.set("w", str(int(plot_area_border_width)))
            solidFill = etree.SubElement(ln, qn("a:solidFill"))
            clr = etree.SubElement(solidFill, qn("a:srgbClr"))
            clr.set("val", f"{plot_area_border}")
            if extLst is not None:
                plot_area_xml.append(extLst)
        except Exception:
            pass

    # Gridlines
    for axis in (chart.value_axis, chart.category_axis):
        if axis is None:
            continue
        try:
            axis.has_major_gridlines = bool(gridlines)
            if gridlines:
                gl = axis.major_gridlines
                gline = gl.format.line
                gline.color.rgb = gridlines_color
                gline.width = gridlines_width
        except Exception:
            pass


def add_bar_chart(
    slide,
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    categories: list[str],
    series_data: list[tuple[str, list[float]]],
    title: str | None = None,
    y_label: str | None = None,
    gridlines: bool = False,
):
    """Excel-editable clustered-column bar chart.

    ``series_data`` is ``[(series_name, values), ...]``. The resulting chart
    shape opens an embedded workbook on double-click so the user can edit
    numbers after the deck is generated.
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_data:
        chart_data.add_series(name, values)
    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    )
    chart = frame.chart
    _style_chart_common(chart, title=title, gridlines=gridlines)
    # カテゴリカル配色（青一色を避ける）
    for i, plot_series in enumerate(chart.series):
        try:
            plot_series.format.fill.solid()
            plot_series.format.fill.fore_color.rgb = CATEGORICAL_RGB[i % len(CATEGORICAL_RGB)]
        except Exception:
            pass
    if y_label is not None:
        try:
            chart.value_axis.has_title = True
            chart.value_axis.axis_title.text_frame.text = y_label
        except Exception:
            pass
    return frame


def add_scatter_line_chart(
    slide,
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    series_data: list[tuple[str, list[float], list[float]]],
    title: str | None = None,
    x_label: str | None = None,
    y_label: str | None = None,
    show_markers: bool = True,
    show_legend: bool = True,
    gridlines: bool = False,
):
    """Excel-editable XY scatter-line chart.

    ``series_data`` is ``[(series_name, x_values, y_values), ...]``.
    """
    from pptx.chart.data import XyChartData
    from pptx.enum.chart import XL_CHART_TYPE

    chart_data = XyChartData()
    for name, xs, ys in series_data:
        s = chart_data.add_series(name)
        for x, y in zip(xs, ys):
            s.add_data_point(x, y)
    chart_type = (
        XL_CHART_TYPE.XY_SCATTER_LINES
        if show_markers
        else XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS
    )
    frame = slide.shapes.add_chart(
        chart_type, left, top, width, height, chart_data
    )
    chart = frame.chart
    _style_chart_common(chart, title=title, show_legend=show_legend,
                        gridlines=gridlines)
    # カテゴリカル配色（線 + マーカーを系列ごとに色分け）
    for i, s in enumerate(chart.series):
        col = CATEGORICAL_RGB[i % len(CATEGORICAL_RGB)]
        try:
            s.format.line.color.rgb = col
        except Exception:
            pass
        # マーカー着色は markers 表示時のみ（NO_MARKERS では no-op + API が raise しうる）
        if show_markers:
            try:
                s.marker.format.fill.solid()
                s.marker.format.fill.fore_color.rgb = col
                s.marker.format.line.color.rgb = col
            except Exception:
                pass
    if x_label is not None:
        try:
            chart.category_axis.has_title = True
            chart.category_axis.axis_title.text_frame.text = x_label
        except Exception:
            pass
    if y_label is not None:
        try:
            chart.value_axis.has_title = True
            chart.value_axis.axis_title.text_frame.text = y_label
        except Exception:
            pass
    return frame


def add_slide_chrome(
    slide, title: str, slide_number: int, total: int | None = None
) -> list[Rect]:
    """Add slide chrome: filled dark-teal title bar, white title, page number.

    参考デザイン: スライドマスター相当の統一見出し。全幅の濃ティールバーに白い
    タイトルとページ番号を載せる。``total`` を渡すと ``N / total`` 形式で表示する。

    Returns the title and page-number rects for ``assert_no_overlap``. The bar
    itself is background chrome and is intentionally excluded from the overlap
    check (the title and number sit on it by design).
    """
    # Filled title bar (background, dark teal, full width)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.78)
    )
    _style_shape_fill(bar, COLOR_BASE_DARK, None)
    # Title text (white, on the bar)
    t_rect = (Inches(0.4), Inches(0.14), Inches(11.0), Inches(0.5))
    add_rich_text_box(
        slide,
        [Paragraph([Run(title, size=Pt(28), bold=True, color=COLOR_WHITE)])],
        left=t_rect[0], top=t_rect[1], width=t_rect[2], height=t_rect[3],
    )
    # Page number "N / total" (white), right-aligned
    num_text = f"{slide_number} / {total}" if total else str(slide_number)
    n_rect = (Inches(11.55), Inches(0.18), Inches(1.35), Inches(0.42))
    add_rich_text_box(
        slide,
        [Paragraph(
            [Run(num_text, size=Pt(14), bold=True, color=COLOR_WHITE, font=FONT_EN)],
            alignment=PP_ALIGN.RIGHT,
        )],
        left=n_rect[0], top=n_rect[1], width=n_rect[2], height=n_rect[3],
    )
    return [
        (*t_rect, "<title>"),
        (*n_rect, "<slide-number>"),
    ]


def add_subheading_bar(
    slide,
    text: str,
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu = Inches(0.36),
) -> Rect:
    """Light-teal filled sub-heading bar (多パネルスライドの各図の見出し).

    参考デザイン: 1 スライドに複数の図を並べるとき、各図の上に淡ティールの小見出し
    バーを置いて何の図かを示す。返り値の rect を ``assert_no_overlap`` に渡す
    (テキストはバー内に収まる前提なので、バー 1 個ぶんの rect だけ返す)。
    """
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _style_shape_fill(bar, COLOR_BASE_LIGHT, None)
    # Label lives in the bar's own text frame (one shape only) so the returned
    # rect fully covers the visible content — same pattern as add_flow_box.
    tf = bar.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    _write_paragraphs(
        tf, [Paragraph(mixed_runs(text, size=Pt(20), bold=True, color=COLOR_TITLE))]
    )
    return (left, top, width, height, f"<subheading:{text[:10]}>")


# ---------------------------------------------------------------------------
# Legacy data-table helper
# ---------------------------------------------------------------------------


def add_data_table(
    slide,
    headers: list[str],
    rows: list[list[str]],
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    header_fill: RGBColor = COLOR_ACCENT_BLUE,
    header_text: RGBColor = RGBColor(0xFF, 0xFF, 0xFF),
    body_text: RGBColor = COLOR_TEXT_BODY,
    font_size: Pt = Pt(16),
    header_size: Pt = Pt(16),
    first_col_bold: bool = True,
    row_height: Emu | None = None,
    cell_margin: Emu = Pt(3),
    word_wrap: bool = False,
    highlight_row: int | None = None,
    highlight_fill: RGBColor = COLOR_ROW_HIGHLIGHT_FILL,
):
    """Add an editable native PowerPoint table.

    Font is auto-selected per cell (MS Gothic for JA, Arial for EN) to honor
    section 1 of the style guide. Header row is filled with ``header_fill``.

    Color (§3 — 青一色を避ける):
    - ``highlight_row``: 1-based body row index to fill with ``highlight_fill``
      (淡いアンバー既定）。比較表で「推奨案 / hero 行」に目を引かせる。
      例: FF/DFT/MLIP 比較で MLIP 行を highlight_row=3。

    Whitespace control (§0 — avoid "AI 作成感"):
    - ``row_height``: explicit per-row height (applies to every row). **Default
      (None / 0) auto-fits**: header row to ``header_size × 2.0`` and body rows
      to ``font_size × 2.0`` (× 2.0 leaves CJK glyph headroom). ``height`` only
      sets the table's initial bounding box; the final shape height is the sum
      of row heights.
    - ``cell_margin``: top/bottom inner margin (default 3 pt).
    - ``word_wrap``: **default False** — data-table cells should hold short
      single-line entries. Wrapping silently expands the tight row height (the
      OOXML row height is a minimum, not a clip), which reintroduces the
      AI-looking whitespace. Set ``True`` only for genuinely long cells, and
      bump ``row_height`` accordingly.

    Rule of thumb: **size the table to its content, not to the slide.** A
    4-row table of short entries is ~1.6" tall. Center a narrow table rather
    than stretching it full-width.
    """
    if highlight_row is not None and not (1 <= highlight_row <= len(rows)):
        raise ValueError(
            f"highlight_row={highlight_row} out of range; use 1-based body row "
            f"index 1..{len(rows)}"
        )
    n_rows = len(rows) + 1
    n_cols = len(headers)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = shape.table

    # Tight per-row heights so cells hug their text (kills AI-looking whitespace).
    # Header and body sized separately so a large header does not inflate body rows.
    # ``Emu(0)`` is treated as "auto" (falsy guard) per review.
    header_rh = int(int(header_size) * 2.0)
    body_rh = int(int(font_size) * 2.0)
    explicit_rh = int(row_height) if row_height else None
    for r in range(n_rows):
        tbl.rows[r].height = explicit_rh if explicit_rh else (header_rh if r == 0 else body_rh)

    def _fill_cell(cell, text, *, size, bold, color, align):
        cell.text_frame.word_wrap = word_wrap
        cell.margin_top = int(cell_margin)
        cell.margin_bottom = int(cell_margin)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = size
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT_JA if _is_japanese(text) else FONT_EN

    for c, htext in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        _fill_cell(cell, htext, size=header_size, bold=True,
                   color=header_text, align=PP_ALIGN.CENTER)

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            if highlight_row is not None and r == highlight_row:
                cell.fill.solid()
                cell.fill.fore_color.rgb = highlight_fill
            _fill_cell(cell, value, size=font_size, bold=(first_col_bold and c == 0),
                       color=body_text, align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)

    return shape


# ---------------------------------------------------------------------------
# §0 absolute-rule enforcement (text minimization)
# ---------------------------------------------------------------------------


def assert_text_minimal(
    slide,
    *,
    max_textboxes: int = 5,
    max_chars_per_box: int = 120,
    max_total_lines: int = 12,
) -> None:
    """Enforce §0 (text minimization). Raise ``ValueError`` on violation.

    Counts shapes with non-empty ``text_frame.text``. Default ``max_textboxes=5``
    accounts for chrome (title + slide_number) + body content (key-message band +
    1 main body shape + 1 supplement) per §0:

      タイトル + 本文 + key-message band = 3 個まで, 補足ラベル 1 個まで許容
      (= 4 content boxes), plus chrome's slide_number = 5 total textboxes.

    Layouts that legitimately need more boxes (e.g. ``split_2col`` which adds
    2 body cards) should pass ``max_textboxes=6`` explicitly.

    Other §0 limits enforced here:
      - 本文ブロックは総 8 行まで（chrome の title + slide_number + key-msg を足して max 12）
      - 1 ボックス内 120 字まで

    Call at the end of each slide builder, alongside :func:`assert_no_overlap`.
    """
    text_shapes: list[tuple[object, str]] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        text_shapes.append((shape, text))

    n_boxes = len(text_shapes)
    if n_boxes > max_textboxes:
        raise ValueError(
            f"§0 violation (textbox count): {n_boxes} text boxes on slide "
            f"(max {max_textboxes} = title + body + key-message + 1 optional). "
            f"Reduce text or split slide."
        )

    total_lines = 0
    for shape, text in text_shapes:
        paras = [p for p in shape.text_frame.paragraphs if p.text.strip()]
        total_lines += len(paras)
        n_chars = len(text)
        if n_chars > max_chars_per_box:
            raise ValueError(
                f"§0 violation (per-box char count): {n_chars} chars (max "
                f"{max_chars_per_box}). First 60 chars: {text[:60]!r}. "
                f"Reduce text or split into multiple shapes."
            )

    if total_lines > max_total_lines:
        raise ValueError(
            f"§0 violation (total lines): {total_lines} non-empty paragraphs "
            f"(max {max_total_lines} = ~8 body + title + key-message + margin). "
            f"Reduce body content or split slide."
        )


# Generic non-assertive title heads to warn against. 結語/Conclusion/Summary are
# intentionally NOT here -- they are valid recap-slide titles. The point of this
# lint is to push data/insight slides toward assertive headlines like
# "Form I が 175 分で Form II に転移" instead of bland "結果".
_TITLE_BLACKLIST_GENERIC: frozenset[str] = frozenset({
    # 日本語: 無味な見出し
    "結果", "考察", "方法", "実験", "実験結果", "解析", "解析結果",
    "目的", "前提", "実験手順", "手順", "概要", "背景", "今後の予定",
    # English: bland section headers
    "results", "discussion", "method", "methods", "experimental",
    "introduction", "background", "next steps", "purpose",
    "objective", "objectives", "procedure", "assumptions", "overview",
})
# Precomputed lowercase set for fast lookup (avoid per-call rebuild)
_TITLE_BLACKLIST_LOWER: frozenset[str] = frozenset(
    b.lower() for b in _TITLE_BLACKLIST_GENERIC
)


def assert_title_assertive(
    title: str,
    *,
    blacklist: frozenset[str] | None = None,
) -> None:
    """Warn (raise ``ValueError``) if a slide title is a generic non-assertive heading.

    Per §0: titles should be assertive claims (e.g., "Form I が 175 分で Form II
    に転移") rather than bland headings ("結果", "考察"). The blacklist is
    case-insensitive and matches the title's trimmed text exactly.

    Override with ``blacklist=frozenset()`` to disable, or pass a custom set.
    """
    if blacklist is None:
        bl_lower = _TITLE_BLACKLIST_LOWER  # precomputed at module load
    else:
        bl_lower = frozenset(b.lower() for b in blacklist)
    normalized = title.strip().lower()
    if normalized in bl_lower:
        raise ValueError(
            f"§0 violation (title not assertive): {title!r}. Use a specific "
            f"claim (e.g. 'Form I が 175 分で Form II に転移') instead of a "
            f"generic heading. Override with blacklist=frozenset() if "
            f"intentional (e.g. for a section divider)."
        )


# ---------------------------------------------------------------------------
# Lazy matplotlib JA configuration (used by helpers that render JA text)
# ---------------------------------------------------------------------------

_matplotlib_ja_configured: bool = False


def _ensure_matplotlib_japanese() -> None:
    """Configure matplotlib to render Japanese once per process.

    Helpers that draw Japanese text via matplotlib (e.g. ``add_timeline``,
    ``add_energy_diagram``) call this on first use. Emits a ``RuntimeWarning``
    and falls back to the matplotlib default (DejaVu Sans, will tofu on CJK)
    if MS Gothic is missing or unreadable. English-only diagrams still work.
    """
    global _matplotlib_ja_configured
    if _matplotlib_ja_configured:
        return
    try:
        configure_matplotlib_japanese()
    except (FileNotFoundError, OSError) as exc:
        warnings.warn(
            f"No Japanese font configured ({exc}). Japanese text in matplotlib "
            f"diagrams will render as tofu (□). Set CAW_SLIDES_MSGOTHIC "
            f"env var to your msgothic.ttc path to fix.",
            RuntimeWarning,
            stacklevel=3,
        )
    _matplotlib_ja_configured = True


# ---------------------------------------------------------------------------
# Chemistry-specific visual helpers
# (RDKit + matplotlib are lazy-imported to keep base import light)
# ---------------------------------------------------------------------------


def add_molecule(
    slide,
    smiles: str,
    *,
    left: Emu,
    top: Emu,
    max_width: Emu,
    max_height: Emu,
    img_size: tuple[int, int] = (400, 400),
):
    """Render a SMILES string via RDKit and embed the resulting PNG.

    The PNG is sized to ``img_size`` pixels then scaled into the bounding box
    via :func:`add_picture_fit` (aspect-preserving, centered).

    Requires ``rdkit`` to be installed (``pip install rdkit``).
    """
    try:
        from rdkit import Chem
        from rdkit.Chem import Draw
    except ImportError as exc:  # pragma: no cover -- import error path
        raise ImportError(
            "rdkit is required for add_molecule. Install with: pip install rdkit"
        ) from exc

    mol = Chem.MolFromSmiles(smiles)
    if mol is None:
        raise ValueError(f"Invalid SMILES: {smiles!r}")

    img = Draw.MolToImage(mol, size=img_size)

    tmp_path: str | None = None
    try:
        # Close-then-reopen pattern: NamedTemporaryFile holds the file open on
        # Windows, which blocks PIL.Image.open inside add_picture_fit. Closing
        # before re-opening makes the helper Windows-safe.
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            tmp_path = tmp.name
        img.save(tmp_path, "PNG")
        return add_picture_fit(
            slide, tmp_path,
            left=left, top=top, max_width=max_width, max_height=max_height,
        )
    finally:
        if tmp_path and os.path.exists(tmp_path):
            os.unlink(tmp_path)


def add_reaction_scheme(
    slide,
    *,
    reactants: list[str],
    products: list[str],
    conditions: str = "",
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
):
    """Compose a reaction scheme horizontally: reactants → products with conditions label.

    ``reactants`` and ``products`` are lists of SMILES strings rendered via
    :func:`add_molecule`. Molecules are joined with "+" between same-side
    entries and a right arrow between the two sides. Conditions text is placed
    above the arrow.

    All molecules share equal horizontal width inside the bounding box.
    """
    n_r, n_p = len(reactants), len(products)
    n_mols = n_r + n_p
    if n_mols == 0:
        raise ValueError("reactants and products cannot both be empty")

    arrow_w = int(Inches(0.8))
    plus_w = int(Inches(0.3))
    n_plus = max(0, n_r - 1) + max(0, n_p - 1)
    mol_w = (int(width) - arrow_w - n_plus * plus_w) // n_mols
    mol_h = int(height)
    if mol_w <= 0:
        raise ValueError(
            f"add_reaction_scheme: width {width} too small for {n_mols} "
            f"molecules + arrow + {n_plus} plus signs"
        )

    def _plus_at(x: int) -> None:
        add_rich_text_box(
            slide,
            [Paragraph(
                [Run("+", size=Pt(28), bold=True, color=COLOR_TEXT_BODY, font=FONT_EN)],
                alignment=PP_ALIGN.CENTER,
            )],
            left=x, top=top, width=plus_w, height=mol_h,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    x = int(left)
    # Reactants
    for i, smi in enumerate(reactants):
        add_molecule(slide, smi, left=x, top=top, max_width=mol_w, max_height=mol_h)
        x += mol_w
        if i < n_r - 1:
            _plus_at(x)
            x += plus_w

    # Arrow (vertically centered) with conditions label above.
    # Clamp arrow_h to mol_h so arrow never extends outside the bounding box,
    # which also keeps the conditions label height non-negative.
    arrow_h = min(int(Inches(0.4)), mol_h)
    arrow_top = int(top) + (mol_h - arrow_h) // 2
    add_flow_arrow(
        slide, left=x, top=arrow_top,
        width=arrow_w, height=arrow_h, direction="right",
    )
    cond_h = arrow_top - int(top)
    if conditions and cond_h > 0:
        add_rich_text_box(
            slide,
            [Paragraph(
                mixed_runs(conditions, size=Pt(12), color=COLOR_TEXT_BODY),
                alignment=PP_ALIGN.CENTER,
            )],
            left=x, top=int(top), width=arrow_w, height=cond_h,
            anchor=MSO_ANCHOR.BOTTOM,
        )
    x += arrow_w

    # Products
    for i, smi in enumerate(products):
        add_molecule(slide, smi, left=x, top=top, max_width=mol_w, max_height=mol_h)
        x += mol_w
        if i < n_p - 1:
            _plus_at(x)
            x += plus_w


def add_energy_diagram(
    slide,
    *,
    levels: list[float],
    labels: list[str],
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    y_label: str = "Energy (kcal/mol)",
):
    """Render a reaction-coordinate energy diagram via matplotlib and embed as PNG.

    Each ``level`` is drawn as a short horizontal segment with the corresponding
    ``label`` above it. Adjacent levels are connected with smooth interpolation
    curves to visualize transitions (TS, intermediates).

    ``levels`` are interpreted in the same units as ``y_label`` implies.
    """
    if len(levels) != len(labels):
        raise ValueError(
            f"levels ({len(levels)}) and labels ({len(labels)}) must match"
        )
    if len(levels) < 2:
        raise ValueError("need at least 2 levels for an energy diagram")

    import numpy as np

    _ensure_matplotlib_japanese()
    n = len(levels)
    fig_w = max(4.0, float(width) / 914400.0)  # EMU -> inches
    fig_h = max(2.0, float(height) / 914400.0)
    fig, ax = plt.subplots(figsize=(fig_w, fig_h), dpi=150)
    connector_rgb = "#888888"  # 状態間をつなぐ破線は中立グレー

    # Endpoint segments are wider (stable reactant/product), intermediate
    # segments (TS, intermediates) are short dashes so they read as a peak
    # rather than a flat plateau.
    def _seg_width(idx: int) -> float:
        return 0.6 if idx in (0, n - 1) else 0.15

    # State role colors: 始状態=teal / 終状態=緑 / 中間(TS・中間体)=テラコッタ
    # （参考デザインのパレット。役割が色で読めると初学者にも障壁/安定性が直感的）
    def _state_color(idx: int) -> str:
        if idx == 0:
            return "#3686A6"   # reactant base teal
        if idx == n - 1:
            return "#459B2D"   # product 緑
        return "#BA5936"       # TS / intermediate テラコッタ

    for i, (e, label) in enumerate(zip(levels, labels)):
        sw = _seg_width(i)
        ax.plot([i - sw / 2, i + sw / 2], [e, e], color=_state_color(i), lw=4)
        ax.annotate(
            label, xy=(i, e), xytext=(0, 8),
            textcoords="offset points", ha="center",
            fontsize=11, fontweight="bold", color=_state_color(i),
        )
        if i < n - 1:
            x0 = i + sw / 2
            x1 = (i + 1) - _seg_width(i + 1) / 2
            xs = np.linspace(x0, x1, 30)
            t = (xs - x0) / (x1 - x0)
            # Smoothstep for a gentler curve than linear
            t_smooth = t * t * (3 - 2 * t)
            curve = e + (levels[i + 1] - e) * t_smooth
            ax.plot(xs, curve, color=connector_rgb, lw=2, ls="--", alpha=0.6)

    ax.set_xlim(-0.5, n - 0.5)
    ax.set_xticks([])
    ax.set_ylabel(y_label, fontsize=12)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["bottom"].set_visible(False)
    fig.tight_layout()

    tmp_path: str | None = None
    try:
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            tmp_path = tmp.name
        fig.savefig(tmp_path, dpi=150, transparent=True)
        plt.close(fig)
        return add_picture_fit(
            slide, tmp_path,
            left=left, top=top, max_width=width, max_height=height,
        )
    finally:
        plt.close(fig)  # idempotent if already closed
        if tmp_path and os.path.exists(tmp_path):
            os.unlink(tmp_path)


# ---------------------------------------------------------------------------
# Layout patterns
# ---------------------------------------------------------------------------


def split_2col(
    slide,
    *,
    left_paragraphs: list[Paragraph],
    right_paragraphs: list[Paragraph],
    top: Emu,
    height: Emu,
    left_title: str | None = None,
    right_title: str | None = None,
    ratio: float = 0.5,
    body_left: Emu = Inches(0.4),
    body_right_margin: Emu = Inches(0.4),
    gap: Emu = Inches(0.2),
):
    """Two-column comparison layout. Returns the (left_shape, right_shape) tuple.

    The left column gets ``ratio`` of the available horizontal space (default
    50/50). Titles, if provided, are added as 21pt bold navy headers above the
    body paragraphs of each column. Both columns use :func:`add_shape_card`
    with default L4 styling (no border / no fill / no shadow).

    Use this whenever §0 says "概念の対比 → 2 カラム" (e.g., Form I vs Form II,
    experiment vs simulation, before vs after).
    """
    if not (0.0 < ratio < 1.0):
        raise ValueError(f"ratio must be in (0, 1), got {ratio}")

    total_avail = int(SLIDE_WIDTH) - int(body_left) - int(body_right_margin) - int(gap)
    if total_avail <= 0:
        raise ValueError(
            f"split_2col: margins + gap exceed slide width "
            f"(body_left={body_left}, body_right_margin={body_right_margin}, "
            f"gap={gap}). Increase slide width or reduce margins."
        )
    left_w = int(total_avail * ratio)
    right_w = total_avail - left_w
    left_x = int(body_left)
    right_x = left_x + left_w + int(gap)

    if left_title:
        left_paragraphs = [
            Paragraph(mixed_runs(left_title, size=Pt(21), bold=True, color=COLOR_TITLE)),
            *left_paragraphs,
        ]
    if right_title:
        right_paragraphs = [
            Paragraph(mixed_runs(right_title, size=Pt(21), bold=True, color=COLOR_TITLE)),
            *right_paragraphs,
        ]

    left_shape = add_shape_card(
        slide, left=left_x, top=top, width=left_w, height=height,
        paragraphs=left_paragraphs,
    )
    right_shape = add_shape_card(
        slide, left=right_x, top=top, width=right_w, height=height,
        paragraphs=right_paragraphs,
    )
    return left_shape, right_shape


def add_timeline(
    slide,
    *,
    milestones: list[tuple[str, str]],
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
):
    """Horizontal timeline: dates as bottom labels, events as top labels.

    ``milestones`` is a list of ``(date_label, event_label)`` tuples, rendered
    left-to-right with equal spacing. Useful for 先行研究年表, 実験スケジュール,
    プロジェクト Milestone in 報告会 / 申請書 slides.
    """
    if not milestones:
        raise ValueError("milestones cannot be empty")
    if len(milestones) > 8:
        raise ValueError(
            f"add_timeline: {len(milestones)} milestones is too dense for one "
            f"slide (labels will collide). Split into multiple slides or "
            f"reduce milestones to ≤ 8."
        )

    _ensure_matplotlib_japanese()
    n = len(milestones)
    fig_w = max(4.0, float(width) / 914400.0)
    fig_h = max(1.2, float(height) / 914400.0)
    fig, ax = plt.subplots(figsize=(fig_w, fig_h), dpi=150)
    line_rgb = "#1E4E63"  # 参考デザイン: 暗ティール（タイムライン軸）

    ax.plot([0.05, 0.95], [0, 0], color=line_rgb, lw=3)
    for i, (date, event) in enumerate(milestones):
        x = 0.05 + 0.9 * (i / max(1, n - 1)) if n > 1 else 0.5
        ax.scatter([x], [0], color=line_rgb, s=200, zorder=3)
        ax.annotate(
            date, xy=(x, 0), xytext=(0, -22),
            textcoords="offset points", ha="center",
            fontsize=10, color=line_rgb, fontweight="bold",
        )
        ax.annotate(
            event, xy=(x, 0), xytext=(0, 16),
            textcoords="offset points", ha="center",
            fontsize=11, color="#222222",
        )

    ax.set_xlim(-0.05, 1.05)
    ax.set_ylim(-1, 1)
    ax.axis("off")
    fig.tight_layout()

    tmp_path: str | None = None
    try:
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            tmp_path = tmp.name
        fig.savefig(tmp_path, dpi=150, transparent=True)
        plt.close(fig)
        return add_picture_fit(
            slide, tmp_path,
            left=left, top=top, max_width=width, max_height=height,
        )
    finally:
        plt.close(fig)
        if tmp_path and os.path.exists(tmp_path):
            os.unlink(tmp_path)
