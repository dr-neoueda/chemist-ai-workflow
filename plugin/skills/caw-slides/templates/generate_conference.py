"""Conference presentation template (oral or poster, 20–25 slides).

caw-slides variant: **学会発表**（国内・国際、口頭・ポスター）

Use case
--------
Specialist audience, fixed time slot (12–15 min oral), lab-color or conference theme.
Slides emphasize: research question → method → results → discussion → conclusion.

Usage
-----
1. Copy this file to ``.company/presentation/scripts/generate_<theme>_<YYYYMMDD>.py``
2. Replace ``<...>`` placeholders with your content
3. Run::

       python .company/presentation/scripts/generate_<theme>_<YYYYMMDD>.py

   The .pptx is written to ``presentations/slides/<theme>_<YYYYMMDD>.pptx``.

Style guide
-----------
Follows ``../references/style-guide.md`` (caw-slides). Uses helpers from
``../references/pptx_helpers.py``: ``add_slide_chrome``, ``add_shape_card``,
``add_key_message_band``, ``assert_no_overlap``, etc.
"""
from __future__ import annotations

import sys
from pathlib import Path

# Locate ``pptx_helpers.py`` in this priority:
#   1. Same directory as this script (after the user copies pptx_helpers.py
#      alongside ``generate_<purpose>_<YYYYMMDD>.py`` in
#      ``.company/presentation/scripts/``)
#   2. ``../references/`` (running this script directly from the plugin install)
# The caw-slides Skill workflow copies ``pptx_helpers.py`` into the user's
# project automatically, so case 1 covers the documented usage.
_HERE = Path(__file__).resolve().parent
for _cand in (_HERE, _HERE.parent / "references"):
    if (_cand / "pptx_helpers.py").is_file():
        _cand_str = str(_cand)
        if _cand_str not in sys.path:
            sys.path.insert(0, _cand_str)
        break

import pptx_helpers as h  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

# ─── Output path ──────────────────────────────────────────────────────────────
# By the caw two-layer principle, .pptx artifacts go to a top-level visible folder.
OUT_DIR = Path.cwd() / "presentations" / "slides"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_PATH = OUT_DIR / "conference_template.pptx"

# ─── Placeholders (replace with your content) ─────────────────────────────────
TITLE_JA = "<研究タイトル>"
PRESENTER = "<your_name>"
COAUTHORS = "<co-author 1>, <co-author 2>"
AFFILIATION = "<affiliation>"
EVENT = "<conference name, YYYY-MM-DD>"


def build_title_slide(prs):
    slide = h.blank_slide(prs)
    rects = []

    # Title (32pt bold, navy)
    h.add_rich_text_box(
        slide,
        [h.Paragraph(h.mixed_runs(TITLE_JA, size=Pt(32), bold=True, color=h.COLOR_TITLE),
                     alignment=PP_ALIGN.CENTER)],
        left=Inches(1.0), top=Inches(2.4), width=Inches(11.3), height=Inches(1.0),
    )
    rects.append((Inches(1.0), Inches(2.4), Inches(11.3), Inches(1.0), "<title>"))

    # Presenter line
    presenter_text = f"○{PRESENTER}, {COAUTHORS} ({AFFILIATION})"
    h.add_rich_text_box(
        slide,
        [h.Paragraph(h.mixed_runs(presenter_text, size=Pt(24), color=h.COLOR_TEXT_BODY),
                     alignment=PP_ALIGN.CENTER)],
        left=Inches(1.0), top=Inches(4.0), width=Inches(11.3), height=Inches(0.6),
    )
    rects.append((Inches(1.0), Inches(4.0), Inches(11.3), Inches(0.6), "<presenter>"))

    # Event line
    h.add_rich_text_box(
        slide,
        [h.Paragraph(h.mixed_runs(EVENT, size=Pt(20), color=h.COLOR_LABEL_GREY),
                     alignment=PP_ALIGN.CENTER)],
        left=Inches(1.0), top=Inches(5.0), width=Inches(11.3), height=Inches(0.5),
    )
    rects.append((Inches(1.0), Inches(5.0), Inches(11.3), Inches(0.5), "<event>"))

    h.assert_no_overlap(rects)


def build_motivation_slide(prs, slide_number: int):
    slide = h.blank_slide(prs)
    rects = list(h.add_slide_chrome(slide, "研究背景・目的", slide_number))

    h.add_shape_card(
        slide,
        left=Inches(0.4), top=Inches(1.18), width=Inches(12.53), height=Inches(4.95),
        paragraphs=[
            h.Paragraph(h.mixed_runs("▸ 先行研究の到達点と残された課題",
                                     size=Pt(21), bold=True, color=h.COLOR_TITLE)),
            h.Paragraph(h.mixed_runs("・<課題 1>", size=Pt(16))),
            h.Paragraph(h.mixed_runs("・<課題 2>", size=Pt(16))),
            h.Paragraph(h.mixed_runs("", size=Pt(16))),
            h.Paragraph(h.mixed_runs("▸ 本研究の目的",
                                     size=Pt(21), bold=True, color=h.COLOR_TITLE)),
            h.Paragraph(h.mixed_runs("・<目的 1>", size=Pt(16))),
            h.Paragraph(h.mixed_runs("・<目的 2>", size=Pt(16))),
        ],
    )
    rects.append((Inches(0.4), Inches(1.18), Inches(12.53), Inches(4.95), "<motivation>"))

    h.add_key_message_band(
        slide,
        [h.Paragraph(h.mixed_runs("<L1 一行: なぜこの研究が必要か>",
                                  size=Pt(20), bold=True, color=h.COLOR_TEXT_BODY),
                     alignment=PP_ALIGN.CENTER)],
    )
    rects.append(h.CODEX_KEY_MSG_RECT + ("<key-msg>",))

    h.add_source_line(slide, "")
    rects.append(h.CODEX_SOURCE_RECT + ("<source>",))
    h.assert_no_overlap(rects)


def build_results_slide(prs, slide_number: int):
    """Conference-typical result slide: 1 chart + key takeaway band."""
    slide = h.blank_slide(prs)
    rects = list(h.add_slide_chrome(slide, "結果: <数値・関係性のサマリ>", slide_number))

    # Native Excel-editable chart (placeholder data)
    h.add_scatter_line_chart(
        slide,
        left=Inches(0.4), top=Inches(1.18), width=Inches(12.53), height=Inches(4.6),
        series_data=[
            ("<series 1>", [0, 1, 2, 3, 4, 5], [1.0, 1.6, 2.1, 2.4, 2.5, 2.6]),
            ("<series 2>", [0, 1, 2, 3, 4, 5], [0.9, 1.2, 1.4, 1.5, 1.5, 1.5]),
        ],
        title=None, x_label="<x-axis label>", y_label="<y-axis label>",
        gridlines=False,
    )
    rects.append((Inches(0.4), Inches(1.18), Inches(12.53), Inches(4.6), "<chart>"))

    h.add_key_message_band(
        slide,
        [h.Paragraph(h.mixed_runs("<L1 一行: 結果が示す主張>",
                                  size=Pt(20), bold=True, color=h.COLOR_TEXT_BODY),
                     alignment=PP_ALIGN.CENTER)],
    )
    rects.append(h.CODEX_KEY_MSG_RECT + ("<key-msg>",))

    h.add_source_line(slide, "")
    rects.append(h.CODEX_SOURCE_RECT + ("<source>",))
    h.assert_no_overlap(rects)


def build_conclusion_slide(prs, slide_number: int):
    slide = h.blank_slide(prs)
    rects = list(h.add_slide_chrome(slide, "結語・今後の展望", slide_number))

    h.add_shape_card(
        slide,
        left=Inches(0.4), top=Inches(1.18), width=Inches(12.53), height=Inches(4.95),
        paragraphs=[
            h.Paragraph(h.mixed_runs("▸ まとめ",
                                     size=Pt(21), bold=True, color=h.COLOR_TITLE)),
            h.Paragraph(h.mixed_runs("○ <主要な結論 1>", size=Pt(16))),
            h.Paragraph(h.mixed_runs("○ <主要な結論 2>", size=Pt(16))),
            h.Paragraph(h.mixed_runs("", size=Pt(16))),
            h.Paragraph(h.mixed_runs("▸ 今後の展望",
                                     size=Pt(21), bold=True, color=h.COLOR_TITLE)),
            h.Paragraph(h.mixed_runs("○ <今後の課題 1>", size=Pt(16))),
            h.Paragraph(h.mixed_runs("○ <今後の課題 2>", size=Pt(16))),
        ],
    )
    rects.append((Inches(0.4), Inches(1.18), Inches(12.53), Inches(4.95), "<conclusion>"))

    h.add_key_message_band(
        slide,
        [h.Paragraph(h.mixed_runs("<L1 一行: この研究で何が変わったか>",
                                  size=Pt(20), bold=True, color=h.COLOR_TEXT_BODY),
                     alignment=PP_ALIGN.CENTER)],
    )
    rects.append(h.CODEX_KEY_MSG_RECT + ("<key-msg>",))

    h.add_source_line(slide, "")
    rects.append(h.CODEX_SOURCE_RECT + ("<source>",))
    h.assert_no_overlap(rects)


def main() -> None:
    prs = h.new_presentation()
    build_title_slide(prs)
    build_motivation_slide(prs, slide_number=2)
    build_results_slide(prs, slide_number=3)
    build_conclusion_slide(prs, slide_number=4)
    prs.save(str(OUT_PATH))
    print(f"Saved {OUT_PATH}")


if __name__ == "__main__":
    main()
