"""Journal club / paper review template (6–12 slides).

caw-slides variant: **論文紹介スライド**（lab 内 journal club、レビュー）

Use case
--------
Introduce a paper to lab colleagues. **Original paper / SI figures are primary**;
self-made figures are supplementary only. Source line must cite the figure number.

Figure extraction workflow (from a PDF placed at ``references/<paper>.pdf``)
---------------------------------------------------------------------------
1. Render each PDF page to PNG at 300 DPI::

       pdftoppm -r 300 references/<paper>.pdf references/page

2. Identify figure pages::

       pdftotext references/<paper>.pdf - | grep -nE "Figure|Scheme|Table"

3. Crop figure regions to ``references/fig_<n>.png`` (PIL.Image.crop or sips -c)
4. Use ``h.add_picture_fit(slide, ...)`` below to occupy 70–90 % of the slide
5. Cite via ``h.add_source_line`` with the figure number

Usage
-----
1. Copy this file to ``office/presentation/scripts/generate_<paper-id>_<YYYYMMDD>.py``
2. Place the paper at ``office/presentation/references/<paper>.pdf``
3. Replace ``<...>`` placeholders and figure paths
4. Run the script — output to ``presentations/slides/<paper-id>_<YYYYMMDD>.pptx``
"""
from __future__ import annotations

import sys
from pathlib import Path

# Locate ``pptx_helpers.py``: script-dir first (after user copies it alongside),
# then ../references/ (running this script directly from the plugin install).
_HERE = Path(__file__).resolve().parent
for _cand in (_HERE, _HERE.parent / "references"):
    if (_cand / "pptx_helpers.py").is_file():
        _cand_str = str(_cand)
        if _cand_str not in sys.path:
            sys.path.insert(0, _cand_str)
        break

import pptx_helpers as h  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

OUT_DIR = Path.cwd() / "presentations" / "slides"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_PATH = OUT_DIR / "journal_club_template.pptx"

# ─── Placeholders ─────────────────────────────────────────────────────────────
PAPER_TITLE = "<Paper title>"
PAPER_AUTHORS = "<First Author> et al."
PAPER_JOURNAL = "<Journal>"
PAPER_YEAR = "<YYYY>"
PAPER_DOI = "<10.xxxx/xxxxx>"
PRESENTER = "<your_name>"


def build_title_slide(prs):
    slide = h.blank_slide(prs)
    rects = []

    h.add_rich_text_box(
        slide,
        [h.Paragraph(h.mixed_runs("論文紹介", size=Pt(28), bold=True, color=h.COLOR_TITLE),
                     alignment=PP_ALIGN.CENTER)],
        left=Inches(1.0), top=Inches(1.6), width=Inches(11.3), height=Inches(0.7),
    )
    rects.append((Inches(1.0), Inches(1.6), Inches(11.3), Inches(0.7), "<label>"))

    h.add_rich_text_box(
        slide,
        [h.Paragraph(h.mixed_runs(PAPER_TITLE, size=Pt(28), bold=True, color=h.COLOR_TEXT_BODY),
                     alignment=PP_ALIGN.CENTER)],
        left=Inches(0.6), top=Inches(2.8), width=Inches(12.1), height=Inches(1.6),
    )
    rects.append((Inches(0.6), Inches(2.8), Inches(12.1), Inches(1.6), "<paper-title>"))

    citation = f"{PAPER_AUTHORS}, {PAPER_JOURNAL} {PAPER_YEAR}  (DOI: {PAPER_DOI})"
    h.add_rich_text_box(
        slide,
        [h.Paragraph(h.mixed_runs(citation, size=Pt(20), color=h.COLOR_TEXT_BODY),
                     alignment=PP_ALIGN.CENTER)],
        left=Inches(0.6), top=Inches(5.0), width=Inches(12.1), height=Inches(0.6),
    )
    rects.append((Inches(0.6), Inches(5.0), Inches(12.1), Inches(0.6), "<citation>"))

    h.add_rich_text_box(
        slide,
        [h.Paragraph(h.mixed_runs(f"presented by {PRESENTER}", size=Pt(20),
                                  color=h.COLOR_LABEL_GREY),
                     alignment=PP_ALIGN.CENTER)],
        left=Inches(0.6), top=Inches(5.8), width=Inches(12.1), height=Inches(0.4),
    )
    rects.append((Inches(0.6), Inches(5.8), Inches(12.1), Inches(0.4), "<presenter>"))

    h.assert_no_overlap(rects)


def build_figure_slide(prs, slide_number: int, title: str,
                       fig_path: Path | None, key_message: str, source: str,
                       support: list[str] | None = None):
    """A slide whose primary content is one original-paper figure PLUS a short
    reading guide.

    ``fig_path`` is a cropped PNG of the figure (700–900 px short edge minimum).
    Set ``fig_path = None`` to render a placeholder rectangle.

    ``support`` = list of 2–4 short strings = **how to read the figure**
    (axes / colours / legend), the key numbers, and the one-line interpretation.
    **Always pass it.** A pasted paper figure with only a single key-message line
    conveys little to the audience — style-guide §3 requires the 3 elements
    図 + L1（key message）+ 支持本文（reading guide）on every content slide,
    journal-club figure slides included.

    Layout: figure on the left, reading-guide card on the right (style-guide §3
    の 2 ゾーン構成). For very wide multi-panel figures, switch to figure-on-top /
    support-card-below instead.
    """
    slide = h.blank_slide(prs)
    rects = list(h.add_slide_chrome(slide, title, slide_number))

    fig_left, fig_top = Inches(0.5), Inches(1.18)
    fig_w, fig_h = Inches(7.7), Inches(5.0)
    card_left, card_w = Inches(8.4), Inches(4.55)

    if fig_path and fig_path.exists():
        h.add_picture_fit(
            slide, fig_path,
            left=fig_left, top=fig_top, max_width=fig_w, max_height=fig_h,
        )
    else:
        # Placeholder rectangle so the layout is overlap-checkable even
        # without the real figure yet extracted.
        h.add_shape_card(
            slide,
            left=fig_left, top=fig_top, width=fig_w, height=fig_h,
            paragraphs=[h.Paragraph(
                h.mixed_runs("[ Figure placeholder — drop the cropped PNG here ]",
                             size=Pt(20), color=h.COLOR_LABEL_GREY),
                alignment=PP_ALIGN.CENTER,
            )],
            border=h.COLOR_CARD_BORDER,
        )
    rects.append((fig_left, fig_top, fig_w, fig_h, "<figure>"))

    # 支持本文（図の読み方・主要数値・解釈）— 図スライドにも必ず添える。
    if support is None:
        support = [
            "<軸・色・凡例が何を表すか>",
            "<どこに注目するか・主要な数値>",
            "<その図が示す結論への含意>",
        ]
    card_paras = [h.Paragraph(h.mixed_runs("▸ 図の読み方", size=Pt(20), bold=True,
                                           color=h.COLOR_TITLE))]
    card_paras += [h.Paragraph(h.mixed_runs("・" + s, size=Pt(20),
                                            color=h.COLOR_TEXT_BODY))
                   for s in support]
    h.add_shape_card(
        slide,
        left=card_left, top=fig_top, width=card_w, height=fig_h,
        paragraphs=card_paras, border=h.COLOR_CARD_BORDER,
    )
    rects.append((card_left, fig_top, card_w, fig_h, "<support>"))

    h.add_key_message_band(
        slide,
        [h.Paragraph(h.mixed_runs(key_message,
                                  size=Pt(24), bold=True, color=h.COLOR_TEXT_BODY),
                     alignment=PP_ALIGN.CENTER)],
    )
    rects.append(h.CODEX_KEY_MSG_RECT + ("<key-msg>",))

    h.add_source_line(slide, source)
    rects.append(h.CODEX_SOURCE_RECT + ("<source>",))
    h.assert_no_overlap(rects)


def build_takeaway_slide(prs, slide_number: int):
    slide = h.blank_slide(prs)
    rects = list(h.add_slide_chrome(slide, "私たちの研究との関係", slide_number))

    h.add_shape_card(
        slide,
        left=Inches(0.4), top=Inches(1.18), width=Inches(12.53), height=Inches(4.95),
        paragraphs=[
            h.Paragraph(h.mixed_runs("▸ この論文から得られる示唆",
                                     size=Pt(24), bold=True, color=h.COLOR_TITLE)),
            h.Paragraph(h.mixed_runs("・<示唆 1>", size=Pt(20))),
            h.Paragraph(h.mixed_runs("・<示唆 2>", size=Pt(20))),
            h.Paragraph(h.mixed_runs("", size=Pt(20))),
            h.Paragraph(h.mixed_runs("▸ 私たちの研究への活かし方",
                                     size=Pt(24), bold=True, color=h.COLOR_TITLE)),
            h.Paragraph(h.mixed_runs("・<活かし方 1>", size=Pt(20))),
            h.Paragraph(h.mixed_runs("・<活かし方 2>", size=Pt(20))),
        ],
    )
    rects.append((Inches(0.4), Inches(1.18), Inches(12.53), Inches(4.95), "<takeaway>"))

    h.add_key_message_band(
        slide,
        [h.Paragraph(h.mixed_runs("<L1 一行: なぜこの論文を選んだのか／何を持ち帰るか>",
                                  size=Pt(24), bold=True, color=h.COLOR_TEXT_BODY),
                     alignment=PP_ALIGN.CENTER)],
    )
    rects.append(h.CODEX_KEY_MSG_RECT + ("<key-msg>",))
    h.assert_no_overlap(rects)


def main() -> None:
    prs = h.new_presentation()
    build_title_slide(prs)
    build_figure_slide(
        prs, slide_number=2, title="<Scheme 1: 反応スキーム>",
        fig_path=Path("office/presentation/references/fig_scheme1.png"),
        key_message="<L1: 反応の特徴を一行で>",
        source=f"Source: {PAPER_AUTHORS}, {PAPER_JOURNAL} {PAPER_YEAR}, Scheme 1.",
    )
    build_figure_slide(
        prs, slide_number=3, title="<Figure 3: 結果プロット>",
        fig_path=Path("office/presentation/references/fig_3.png"),
        key_message="<L1: 結果が示す主張を一行で>",
        source=f"Source: {PAPER_AUTHORS}, {PAPER_JOURNAL} {PAPER_YEAR}, Figure 3.",
    )
    build_takeaway_slide(prs, slide_number=4)
    prs.save(str(OUT_PATH))
    print(f"Saved {OUT_PATH}")


if __name__ == "__main__":
    main()
