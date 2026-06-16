"""Lecture / tutorial template (15–30 slides).

caw-slides variant: **講義・チュートリアル資料**

Use case
--------
Teach a topic to students or researchers in adjacent fields. Plain language,
concept diagrams, and step-by-step progression are emphasized. Less specialist
jargon than a conference deck.

Structure
---------
1. Title slide
2. Learning goals (what the audience will be able to do after the lecture)
3. Topic introduction (1–3 slides, plain-language motivation)
4. Concept slides (5–15 slides, one concept per slide, diagram-first)
5. Hands-on / examples (2–5 slides)
6. Summary / further reading (1 slide)

Usage
-----
1. Copy this file to ``office/presentation/scripts/generate_lecture_<topic>_<YYYYMMDD>.py``
2. Replace placeholders + add concept diagrams
3. Run — output to ``presentations/slides/lecture_<topic>_<YYYYMMDD>.pptx``
"""
from __future__ import annotations

import sys
from pathlib import Path

# Locate ``pptx_helpers.py``: script-dir first (after user copies it alongside),
# then ../references/ (running this script directly from the plugin install).
_HERE = Path(__file__).resolve().parent
for _cand in (_HERE, _HERE.parent / "references"):
    if (_cand / "pptx_helpers.py").is_file():
        _cand_str = str(_cand)
        if _cand_str not in sys.path:
            sys.path.insert(0, _cand_str)
        break

import pptx_helpers as h  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

OUT_DIR = Path.cwd() / "presentations" / "slides"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_PATH = OUT_DIR / "lecture_template.pptx"

# ─── Placeholders ─────────────────────────────────────────────────────────────
LECTURE_TITLE = "<講義タイトル>"
TOPIC = "<トピック>"
INSTRUCTOR = "<your_name>"
AUDIENCE = "<対象聴衆（例: 学部 3 年生、他分野研究者）>"


def build_title_slide(prs):
    slide = h.blank_slide(prs)
    rects = []

    h.add_rich_text_box(
        slide,
        [h.Paragraph(h.mixed_runs(LECTURE_TITLE, size=Pt(32), bold=True, color=h.COLOR_TITLE),
                     alignment=PP_ALIGN.CENTER)],
        left=Inches(1.0), top=Inches(2.4), width=Inches(11.3), height=Inches(1.0),
    )
    rects.append((Inches(1.0), Inches(2.4), Inches(11.3), Inches(1.0), "<title>"))

    h.add_rich_text_box(
        slide,
        [h.Paragraph(h.mixed_runs(f"対象: {AUDIENCE}", size=Pt(20), color=h.COLOR_TEXT_BODY),
                     alignment=PP_ALIGN.CENTER)],
        left=Inches(1.0), top=Inches(4.2), width=Inches(11.3), height=Inches(0.5),
    )
    rects.append((Inches(1.0), Inches(4.2), Inches(11.3), Inches(0.5), "<audience>"))

    h.add_rich_text_box(
        slide,
        [h.Paragraph(h.mixed_runs(INSTRUCTOR, size=Pt(20), color=h.COLOR_LABEL_GREY),
                     alignment=PP_ALIGN.CENTER)],
        left=Inches(1.0), top=Inches(5.0), width=Inches(11.3), height=Inches(0.5),
    )
    rects.append((Inches(1.0), Inches(5.0), Inches(11.3), Inches(0.5), "<instructor>"))
    h.assert_no_overlap(rects)


def build_goals_slide(prs, slide_number: int):
    slide = h.blank_slide(prs)
    rects = list(h.add_slide_chrome(slide, "この講義で身につくこと", slide_number))

    h.add_shape_card(
        slide,
        left=Inches(0.4), top=Inches(1.18), width=Inches(12.53), height=Inches(4.95),
        paragraphs=[
            h.Paragraph(h.mixed_runs("▸ 目標", size=Pt(24), bold=True, color=h.COLOR_TITLE)),
            h.Paragraph(h.mixed_runs("○ <目標 1: 〜できるようになる>", size=Pt(20))),
            h.Paragraph(h.mixed_runs("○ <目標 2: 〜を説明できるようになる>", size=Pt(20))),
            h.Paragraph(h.mixed_runs("○ <目標 3: 〜を計算できるようになる>", size=Pt(20))),
            h.Paragraph(h.mixed_runs("", size=Pt(20))),
            h.Paragraph(h.mixed_runs("▸ 前提知識", size=Pt(24), bold=True, color=h.COLOR_TITLE)),
            h.Paragraph(h.mixed_runs("・<前提 1>", size=Pt(20))),
            h.Paragraph(h.mixed_runs("・<前提 2>", size=Pt(20))),
        ],
    )
    rects.append((Inches(0.4), Inches(1.18), Inches(12.53), Inches(4.95), "<goals>"))

    h.add_key_message_band(
        slide,
        [h.Paragraph(h.mixed_runs("<L1: この講義で一番伝えたいこと>",
                                  size=Pt(24), bold=True, color=h.COLOR_TEXT_BODY),
                     alignment=PP_ALIGN.CENTER)],
    )
    rects.append(h.CODEX_KEY_MSG_RECT + ("<key-msg>",))
    h.assert_no_overlap(rects)


def build_concept_flow_slide(prs, slide_number: int):
    """Concept slide using a native flow diagram (boxes + arrows)."""
    slide = h.blank_slide(prs)
    rects = list(h.add_slide_chrome(slide, "<概念名>: ステップで理解する", slide_number))

    # Three flow boxes + two arrows. Centering math:
    #   walk = box + gap + arrow + gap + box + gap + arrow + gap + box
    #        = 3*box_w + 2*arrow_w + 4*gap
    n_boxes, n_arrows, n_gaps = 3, 2, 4
    box_w, box_h, gap = int(Inches(3.0)), int(Inches(1.4)), int(Inches(0.3))
    arrow_w, arrow_h = int(Inches(0.5)), int(Inches(0.6))
    row_y = int(Inches(3.0))
    total = n_boxes * box_w + n_arrows * arrow_w + n_gaps * gap
    start_x = (int(Inches(13.33)) - total) // 2

    x = start_x
    for i, label in enumerate(["<段階 1>", "<段階 2>", "<段階 3>"]):
        h.add_flow_box(
            slide, label,
            left=x, top=row_y, width=box_w, height=box_h,
            fill=h.COLOR_PILL_BLUE_FILL, text_color=h.COLOR_TEXT_BODY,
            font_size=Pt(18),
        )
        rects.append((x, row_y, box_w, box_h, f"<flow-{i+1}>"))
        x += box_w
        if i < 2:
            x += gap
            h.add_flow_arrow(
                slide,
                left=x, top=row_y + (box_h - arrow_h) // 2,
                width=arrow_w, height=arrow_h, direction="right",
            )
            rects.append((x, row_y + (box_h - arrow_h) // 2, arrow_w, arrow_h, f"<arrow-{i+1}>"))
            x += arrow_w + gap

    h.add_key_message_band(
        slide,
        [h.Paragraph(h.mixed_runs("<L1: 3 段階のうち一番重要な気づき>",
                                  size=Pt(24), bold=True, color=h.COLOR_TEXT_BODY),
                     alignment=PP_ALIGN.CENTER)],
    )
    rects.append(h.CODEX_KEY_MSG_RECT + ("<key-msg>",))
    h.assert_no_overlap(rects)


def build_summary_slide(prs, slide_number: int):
    slide = h.blank_slide(prs)
    rects = list(h.add_slide_chrome(slide, "まとめと次の学習", slide_number))

    h.add_shape_card(
        slide,
        left=Inches(0.4), top=Inches(1.18), width=Inches(12.53), height=Inches(4.95),
        paragraphs=[
            h.Paragraph(h.mixed_runs("▸ 今日のポイント",
                                     size=Pt(24), bold=True, color=h.COLOR_TITLE)),
            h.Paragraph(h.mixed_runs("○ <ポイント 1>", size=Pt(20))),
            h.Paragraph(h.mixed_runs("○ <ポイント 2>", size=Pt(20))),
            h.Paragraph(h.mixed_runs("", size=Pt(20))),
            h.Paragraph(h.mixed_runs("▸ さらに学びたい人へ",
                                     size=Pt(24), bold=True, color=h.COLOR_TITLE)),
            h.Paragraph(h.mixed_runs("・<参考書・論文 1>", size=Pt(20))),
            h.Paragraph(h.mixed_runs("・<参考書・論文 2>", size=Pt(20))),
        ],
    )
    rects.append((Inches(0.4), Inches(1.18), Inches(12.53), Inches(4.95), "<summary>"))

    h.add_key_message_band(
        slide,
        [h.Paragraph(h.mixed_runs(f"<L1: {TOPIC} を学ぶ意義を一行で>",
                                  size=Pt(24), bold=True, color=h.COLOR_TEXT_BODY),
                     alignment=PP_ALIGN.CENTER)],
    )
    rects.append(h.CODEX_KEY_MSG_RECT + ("<key-msg>",))
    h.assert_no_overlap(rects)


def main() -> None:
    prs = h.new_presentation()
    build_title_slide(prs)
    build_goals_slide(prs, slide_number=2)
    build_concept_flow_slide(prs, slide_number=3)
    build_summary_slide(prs, slide_number=4)
    prs.save(str(OUT_PATH))
    print(f"Saved {OUT_PATH}")


if __name__ == "__main__":
    main()
