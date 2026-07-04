"""fix_ea_font のテスト（ビルド後 pptx の CJK run を ea=MS Gothic に強制）。"""
import pytest

import fix_ea_font as fef


@pytest.mark.parametrize(
    "text,expect",
    [
        ("日本語", True),
        ("あいう", True),
        ("① 広帯域", True),   # 丸数字 U+2460（囲み CJK）
        ("Å 記号", True),      # 全角混在
        ("ASCII only 123", False),
        ("kJ/mol", False),
        ("", False),
    ],
)
def test_has_cjk(text, expect):
    assert fef.has_cjk(text) is expect


def _ea_typeface(run):
    from pptx.oxml.ns import qn

    rPr = run._r.find(qn("a:rPr"))
    if rPr is None:
        return None
    ea = rPr.find(qn("a:ea"))
    return None if ea is None else ea.get("typeface")


def _build_pptx(path):
    """CJK run(ea=Arial) と Latin run(ea=Arial) を持つ 1 枚デッキを作る（変換器の癖を模擬）。"""
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    para = tb.text_frame.paragraphs[0]
    r_cjk = para.add_run()
    r_cjk.text = "日本語テキスト"
    fef.get_or_add_ea(r_cjk._r.get_or_add_rPr()).set("typeface", "Arial")  # 癖: ea=Arial
    r_latin = para.add_run()
    r_latin.text = "ASCII"
    fef.get_or_add_ea(r_latin._r.get_or_add_rPr()).set("typeface", "Arial")
    prs.save(str(path))
    return str(path)


def test_fix_forces_ea_on_cjk_runs(tmp_path):
    pptx = pytest.importorskip("pptx")
    path = _build_pptx(tmp_path / "deck.pptx")

    fixed = fef.fix_ea_font(path, ja_font="MS Gothic")
    assert fixed == 1  # CJK run のみ修正

    prs = pptx.Presentation(path)
    runs = list(prs.slides[0].shapes[0].text_frame.paragraphs[0].runs)
    cjk_run = next(r for r in runs if fef.has_cjk(r.text))
    latin_run = next(r for r in runs if not fef.has_cjk(r.text))
    assert _ea_typeface(cjk_run) == "MS Gothic"
    assert _ea_typeface(latin_run) == "Arial"  # Latin run は据え置き


def test_fix_is_idempotent(tmp_path):
    pytest.importorskip("pptx")
    path = _build_pptx(tmp_path / "deck.pptx")
    assert fef.fix_ea_font(path) == 1
    assert fef.fix_ea_font(path) == 0  # 2 回目は変更なし


def test_fix_custom_font(tmp_path):
    pptx = pytest.importorskip("pptx")
    path = _build_pptx(tmp_path / "deck.pptx")
    fef.fix_ea_font(path, ja_font="Yu Gothic")
    prs = pptx.Presentation(path)
    cjk_run = next(
        r for r in prs.slides[0].shapes[0].text_frame.paragraphs[0].runs if fef.has_cjk(r.text)
    )
    assert _ea_typeface(cjk_run) == "Yu Gothic"


def test_fix_recurses_into_groups(tmp_path):
    pptx = pytest.importorskip("pptx")
    from pptx.util import Emu

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    tb = group.shapes.add_textbox(Emu(914400), Emu(914400), Emu(3000000), Emu(900000))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "群の中の日本語"
    fef.get_or_add_ea(run._r.get_or_add_rPr()).set("typeface", "Arial")
    path = str(tmp_path / "grp.pptx")
    prs.save(path)

    assert fef.fix_ea_font(path) == 1


def test_ea_inserted_before_cs():
    """<a:ea> は <a:latin> の後・<a:cs> の前に挿入される（OOXML スキーマ順）。"""
    pytest.importorskip("pptx")
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls

    rPr = parse_xml(
        f'<a:rPr {nsdecls("a")}><a:latin typeface="Arial"/><a:cs typeface="Arial"/></a:rPr>'
    )
    fef.get_or_add_ea(rPr).set("typeface", "MS Gothic")
    order = [t.tag.split("}")[-1] for t in rPr]
    assert order == ["latin", "ea", "cs"]


def test_fix_recurses_into_table_cells(tmp_path):
    """native table の cell 内 CJK run も修正対象（Codex HIGH 回帰）。"""
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table = slide.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(3), Inches(1)).table
    cell = table.cell(0, 0)
    cell.text = "表内の日本語"
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            fef.get_or_add_ea(run._r.get_or_add_rPr()).set("typeface", "Arial")
    path = str(tmp_path / "table.pptx")
    prs.save(path)

    assert fef.fix_ea_font(path) == 1
